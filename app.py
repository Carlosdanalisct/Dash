import argparse
import csv
import io
import hashlib
import importlib.util
import json
import re
import sqlite3
import sys
import threading
import time
import unicodedata
from datetime import datetime
from email.parser import BytesParser
from email.policy import default
from http import HTTPStatus
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from urllib.parse import parse_qs, urlparse

REQUIRED_DEPENDENCIES = {
    "pandas": "pandas",
    "numpy": "numpy",
    "openpyxl": "openpyxl",
    "sqlalchemy": "sqlalchemy",
    "cachetools": "cachetools",
    "apscheduler": "apscheduler",
    "python-pptx": "pptx",
    "reportlab": "reportlab",
    "jinja2": "jinja2",
}


def missing_dependencies():
    return [
        package
        for package, module_name in REQUIRED_DEPENDENCIES.items()
        if importlib.util.find_spec(module_name) is None
    ]


def validate_dependencies(exit_on_missing=False):
    missing = missing_dependencies()
    if not missing:
        return True
    message = (
        "\nDependencias ausentes para iniciar o Dashboard MMS:\n"
        f"- {', '.join(missing)}\n\n"
        "Para corrigir, abra o terminal na pasta do dashboard e execute:\n"
        "pip install -r requirements.txt\n"
    )
    print(message, file=sys.stderr)
    if exit_on_missing:
        raise SystemExit(1)
    return False


if __name__ == "__main__":
    validate_dependencies(exit_on_missing=True)

import pandas as pd


APP_DIR = Path(__file__).resolve().parent
STATIC_DIR = APP_DIR / "static"
UPLOAD_DIR = APP_DIR / "uploads"
DB_PATH = APP_DIR / "reclamacoes.db"
CITIES_PATH = STATIC_DIR / "data" / "cidades_brasil.json"
CONFIG_PATH = APP_DIR / "config.json"
UNKNOWN_UF_VALUES = {"", "NI", "NA", "N/I", "SEM_UF", "SEM UF", "SEM ESTADO", "NAN", "NONE", "NULL"}
GEO_SYNC_LOCK = threading.RLock()
GEO_SYNC_LAST_RUN = 0
GEO_SYNC_TTL_SECONDS = 300

DEFAULT_CONFIG = {
    "daily_closed_goal": 10,
    "sla_business_days_target": 5,
    "date_lag_warning_days": 45,
    "custo_cancelamento": 0,
    "custo_frustracao": 0,
    "custo_improdutiva": 0,
    "limite_alerta_amarelo": 40,
    "limite_alerta_vermelho": 70,
    "auth_enabled": False,
    "efficiency_score": {
        "closed_weight": 1.0,
        "sla_weight": 1.0,
        "backlog_penalty": 0.75,
    },
    "risk_score": {
        "volume_weight": 0.40,
        "procedencia_weight": 0.30,
        "sla_weight": 0.20,
        "reincidencia_weight": 0.10,
    },
}


def load_config():
    config = json.loads(json.dumps(DEFAULT_CONFIG))
    if CONFIG_PATH.exists():
        try:
            loaded = json.loads(CONFIG_PATH.read_text(encoding="utf-8"))
            for key, value in loaded.items():
                if isinstance(value, dict) and isinstance(config.get(key), dict):
                    config[key].update(value)
                else:
                    config[key] = value
        except Exception:
            pass
    return config


def save_config(config):
    CONFIG_PATH.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")


def refresh_config():
    global CONFIG, DAILY_CLOSED_GOAL, SLA_BUSINESS_DAYS_TARGET, DATE_LAG_WARNING_DAYS, EFFICIENCY_SCORE_CONFIG
    CONFIG = load_config()
    DAILY_CLOSED_GOAL = int(CONFIG.get("daily_closed_goal", DEFAULT_CONFIG["daily_closed_goal"]))
    SLA_BUSINESS_DAYS_TARGET = int(CONFIG.get("sla_business_days_target", DEFAULT_CONFIG["sla_business_days_target"]))
    DATE_LAG_WARNING_DAYS = int(CONFIG.get("date_lag_warning_days", DEFAULT_CONFIG["date_lag_warning_days"]))
    EFFICIENCY_SCORE_CONFIG = CONFIG.get("efficiency_score", DEFAULT_CONFIG["efficiency_score"])
    try:
        import auth
        auth.AUTH_ENABLED = bool(CONFIG.get("auth_enabled", False))
    except Exception:
        pass
    return CONFIG


CONFIG = load_config()
try:
    import auth
    auth.AUTH_ENABLED = bool(CONFIG.get("auth_enabled", False))
except Exception:
    pass
DAILY_CLOSED_GOAL = int(CONFIG.get("daily_closed_goal", DEFAULT_CONFIG["daily_closed_goal"]))
SLA_BUSINESS_DAYS_TARGET = int(CONFIG.get("sla_business_days_target", DEFAULT_CONFIG["sla_business_days_target"]))
DATE_LAG_WARNING_DAYS = int(CONFIG.get("date_lag_warning_days", DEFAULT_CONFIG["date_lag_warning_days"]))
EFFICIENCY_SCORE_CONFIG = CONFIG.get("efficiency_score", DEFAULT_CONFIG["efficiency_score"])


def slug(value):
    text = unicodedata.normalize("NFKD", str(value)).encode("ascii", "ignore").decode("ascii")
    return re.sub(r"[^a-z0-9]+", "_", text.lower()).strip("_")


def normalizar_cidade(valor):
    text = unicodedata.normalize("NFKD", str(valor or "")).encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^A-Za-z0-9 ]+", " ", text.upper())
    return re.sub(r"\s+", " ", text).strip()


def text_key(value):
    return slug(normalizar_cidade(value))


def normalize_uf(value):
    text = clean(value, "NI").strip().upper()
    text = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"\s+", " ", text).strip()
    if text in UNKNOWN_UF_VALUES:
        return "NI"
    return text[:2] if len(text) == 2 else text


def is_unknown_uf(value):
    return normalize_uf(value) == "NI"


def find_col(columns, *needles):
    for col in columns:
        normalized = slug(col)
        if all(needle in normalized for needle in needles):
            return col
    raise KeyError(f"Campo nao encontrado: {needles}")


def find_optional_col(columns, *needles):
    try:
        return find_col(columns, *needles)
    except KeyError:
        return None


def clean(value, fallback):
    if pd.isna(value):
        return fallback
    text = str(value).strip()
    if not text or text.lower() == "nan":
        return fallback
    return text


def clean_date(value):
    if pd.isna(value):
        return None
    date = pd.to_datetime(value, errors="coerce")
    if pd.isna(date):
        return None
    return date.strftime("%Y-%m-%d")


def business_days_between(start, end):
    if pd.isna(start) or pd.isna(end):
        return None
    start = pd.to_datetime(start, errors="coerce")
    end = pd.to_datetime(end, errors="coerce")
    if pd.isna(start) or pd.isna(end) or end.date() < start.date():
        return None
    return len(pd.bdate_range(start=start.date(), end=end.date()))


def pct_value(part, total):
    return part / total if total else 0


def status_flag(value, *needles):
    normalized = slug(value)
    return any(needle in normalized for needle in needles)


def add_period_filter(where, args, field, month_field, params, max_key):
    period = params.get("period", "all")
    if period == "last7":
        max_date = params.get(max_key)
        if max_date:
            where.append(f"{field} >= date(?, '-7 day')")
            args.append(max_date)
    elif period == "last30":
        max_date = params.get(max_key)
        if max_date:
            where.append(f"{field} >= date(?, '-30 day')")
            args.append(max_date)
    elif period == "last90":
        max_date = params.get(max_key)
        if max_date:
            where.append(f"{field} >= date(?, '-90 day')")
            args.append(max_date)
    elif period == "last12":
        max_date = params.get(max_key)
        if max_date:
            where.append(f"{field} >= date(?, '-12 month')")
            args.append(max_date)
    elif period == "ytd":
        max_date = params.get(max_key)
        if max_date:
            where.append(f"{field} >= date(?, 'start of year')")
            args.append(max_date)
    elif period == "month" and params.get("month"):
        where.append(f"{month_field} = ?")
        args.append(params["month"])
    elif period == "range":
        if params.get("start"):
            where.append(f"{field} >= ?")
            args.append(params["start"])
        if params.get("end"):
            where.append(f"{field} <= ?")
            args.append(params["end"])


def connect():
    conn = sqlite3.connect(DB_PATH, timeout=30)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA busy_timeout = 30000")
    return conn


def init_db():
    APP_DIR.mkdir(parents=True, exist_ok=True)
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    with connect() as conn:
        conn.executescript(
            """
            CREATE TABLE IF NOT EXISTS import_batches (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                file_name TEXT NOT NULL,
                imported_at TEXT NOT NULL,
                mode TEXT NOT NULL,
                rows_total INTEGER NOT NULL,
                rows_inserted INTEGER NOT NULL,
                procedentes INTEGER NOT NULL,
                date_min TEXT,
                date_max TEXT
            );

            CREATE TABLE IF NOT EXISTS complaints (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                record_hash TEXT NOT NULL,
                data_reclamacao TEXT,
                dt_fechamento TEXT,
                mes TEXT,
                mes_fechamento TEXT,
                cliente TEXT,
                cidade TEXT,
                uf TEXT,
                regiao TEXT,
                hub TEXT,
                servico TEXT,
                status_reclamacao TEXT,
                validacao TEXT,
                motivo TEXT,
                grupo_motivo TEXT,
                prestador TEXT,
                analista_sac TEXT,
                dias_uteis_fechamento INTEGER,
                sla_ok INTEGER,
                procedente INTEGER NOT NULL,
                source_file TEXT NOT NULL,
                imported_at TEXT NOT NULL
            );

            CREATE INDEX IF NOT EXISTS idx_complaints_date ON complaints(data_reclamacao);
            CREATE INDEX IF NOT EXISTS idx_complaints_month ON complaints(mes);
            CREATE INDEX IF NOT EXISTS idx_complaints_proc ON complaints(procedente);
            CREATE INDEX IF NOT EXISTS idx_complaints_dims ON complaints(regiao, hub, cliente, cidade);
            CREATE INDEX IF NOT EXISTS idx_complaints_analyst ON complaints(analista_sac);

            CREATE TABLE IF NOT EXISTS operation_batches (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                file_name TEXT NOT NULL,
                imported_at TEXT NOT NULL,
                mode TEXT NOT NULL,
                rows_total INTEGER NOT NULL,
                rows_inserted INTEGER NOT NULL,
                concluidas INTEGER NOT NULL,
                date_min TEXT,
                date_max TEXT
            );

            CREATE TABLE IF NOT EXISTS operations (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                record_hash TEXT NOT NULL,
                data TEXT,
                mes TEXT,
                numero_assistencia TEXT,
                filial TEXT,
                categoria TEXT,
                cidade TEXT,
                estado TEXT,
                regional TEXT,
                area_trabalho TEXT,
                hub TEXT,
                posto TEXT,
                tipo_atividade TEXT,
                status TEXT,
                recurso TEXT,
                nome_cliente TEXT,
                cliente TEXT,
                motivo_nao_realizada TEXT,
                id_montador TEXT,
                concluida INTEGER NOT NULL,
                cancelada INTEGER NOT NULL,
                frustrada INTEGER NOT NULL,
                reagendada INTEGER NOT NULL,
                improdutiva INTEGER NOT NULL,
                perda INTEGER NOT NULL,
                source_file TEXT NOT NULL,
                imported_at TEXT NOT NULL
            );

            CREATE INDEX IF NOT EXISTS idx_operations_date ON operations(data);
            CREATE INDEX IF NOT EXISTS idx_operations_month ON operations(mes);
            CREATE INDEX IF NOT EXISTS idx_operations_dims ON operations(regional, hub, cidade, cliente);
            CREATE INDEX IF NOT EXISTS idx_operations_resource ON operations(recurso);

            CREATE TABLE IF NOT EXISTS users (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                nome TEXT NOT NULL,
                email TEXT NOT NULL UNIQUE,
                senha_hash TEXT NOT NULL,
                perfil TEXT NOT NULL,
                regional TEXT,
                analista_sac TEXT,
                ativo INTEGER NOT NULL DEFAULT 1,
                created_at TEXT NOT NULL
            );

            CREATE INDEX IF NOT EXISTS idx_users_email ON users(email);
            CREATE INDEX IF NOT EXISTS idx_users_perfil ON users(perfil);

            CREATE TABLE IF NOT EXISTS user_sessions (
                token TEXT PRIMARY KEY,
                user_id INTEGER NOT NULL,
                created_at TEXT NOT NULL,
                expires_at TEXT NOT NULL,
                FOREIGN KEY(user_id) REFERENCES users(id)
            );

            CREATE INDEX IF NOT EXISTS idx_user_sessions_user ON user_sessions(user_id);
            CREATE INDEX IF NOT EXISTS idx_user_sessions_expiry ON user_sessions(expires_at);

            CREATE TABLE IF NOT EXISTS geocidades (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                cidade TEXT NOT NULL,
                uf TEXT NOT NULL,
                cidade_key TEXT NOT NULL,
                uf_key TEXT NOT NULL,
                lat REAL,
                lon REAL,
                fonte TEXT NOT NULL DEFAULT 'base',
                precisa_coordenada INTEGER NOT NULL DEFAULT 1,
                uf_corrigida_de TEXT,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                UNIQUE(cidade_key, uf_key)
            );

            CREATE INDEX IF NOT EXISTS idx_geocidades_keys ON geocidades(cidade_key, uf_key);
            CREATE INDEX IF NOT EXISTS idx_geocidades_missing ON geocidades(precisa_coordenada);

            CREATE TABLE IF NOT EXISTS geocidades_pendentes (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                cidade TEXT NOT NULL,
                uf TEXT,
                uf_original TEXT,
                cidade_normalizada TEXT NOT NULL,
                ocorrencias INTEGER,
                origem TEXT,
                reclamacoes INTEGER NOT NULL DEFAULT 0,
                procedentes INTEGER NOT NULL DEFAULT 0,
                assistencias INTEGER NOT NULL DEFAULT 0,
                ignorada INTEGER NOT NULL DEFAULT 0,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                UNIQUE(cidade_normalizada, uf, origem)
            );

            CREATE INDEX IF NOT EXISTS idx_geocidades_pendentes_lookup
            ON geocidades_pendentes(cidade_normalizada, uf, ignorada);
            """
        )
        if conn.execute("SELECT COUNT(*) AS total FROM users").fetchone()["total"] == 0:
            from auth import create_user
            create_user(conn, "Administrador MMS", "admin@mms.local", "admin123", perfil="admin")
        cols = [row["name"] for row in conn.execute("PRAGMA table_info(complaints)").fetchall()]
        for name, definition in {
            "analista_sac": "TEXT",
            "dt_fechamento": "TEXT",
            "mes_fechamento": "TEXT",
            "dias_uteis_fechamento": "INTEGER",
            "sla_ok": "INTEGER",
        }.items():
            if name not in cols:
                conn.execute(f"ALTER TABLE complaints ADD COLUMN {name} {definition}")


def row_hash(row):
    parts = [
        row.get("id_feedback", ""),
        row.get("id_assistservico", ""),
        row.get("id_assistencia", ""),
        row.get("date", ""),
        row.get("cliente", ""),
        row.get("cidade", ""),
        row.get("hub", ""),
        row.get("motivo", ""),
        row.get("row_index", ""),
    ]
    return hashlib.sha256("||".join(str(part) for part in parts).encode("utf-8")).hexdigest()


def operation_row_hash(row):
    parts = [
        row.get("numero_assistencia", ""),
        row.get("date", ""),
        row.get("cliente", ""),
        row.get("cidade", ""),
        row.get("hub", ""),
        row.get("recurso", ""),
        row.get("status", ""),
        row.get("motivo", ""),
        row.get("row_index", ""),
    ]
    return hashlib.sha256("||".join(str(part) for part in parts).encode("utf-8")).hexdigest()


def excel_sheet_candidates(path):
    candidates = []
    with pd.ExcelFile(path) as workbook:
        for sheet_name in workbook.sheet_names:
            preview = pd.read_excel(workbook, sheet_name=sheet_name, nrows=10)
            columns = {slug(col) for col in preview.columns}
            operation_score = sum(
                1
                for expected in [
                    "numero_da_assistencia",
                    "tipo_atividade",
                    "recurso",
                    "status",
                    "data",
                    "motivo_nao_realizada",
                ]
                if expected in columns
            )
            complaint_score = 0
            for col in columns:
                if col in {"data_reclamacao", "dt_fechamento", "analistasac", "hub"}:
                    complaint_score += 2
                if "reclamacao" in col or "procedencia" in col or "feedback" in col:
                    complaint_score += 1
            sheet_type = "operation" if operation_score >= 3 and operation_score >= complaint_score else "complaints"
            score = max(operation_score, complaint_score)
            candidates.append(
                {
                    "sheet": sheet_name,
                    "type": sheet_type,
                    "score": score,
                    "operationScore": operation_score,
                    "complaintScore": complaint_score,
                }
            )
    candidates.sort(key=lambda item: item["score"], reverse=True)
    return candidates


def read_best_excel_sheet(path):
    candidates = excel_sheet_candidates(path)
    if not candidates or candidates[0]["score"] <= 0:
        sheets = ", ".join(item["sheet"] for item in candidates) or "nenhuma aba encontrada"
        raise ValueError(f"Nao encontrei uma aba compativel para importar. Abas disponiveis: {sheets}.")
    best = candidates[0]
    df = pd.read_excel(path, sheet_name=best["sheet"])
    return df, best


def operation_flags(status, reason):
    combined = f"{status} {reason}"
    concluida = 1 if status_flag(status, "concluida", "concluido", "finalizada", "realizada") else 0
    cancelada = 1 if status_flag(combined, "cancelada", "cancelado", "cancelamento") else 0
    frustrada = 1 if status_flag(combined, "frustrada", "frustrado", "frustracao") else 0
    reagendada = 1 if status_flag(combined, "reagendada", "reagendado", "reagendamento") else 0
    improdutiva = 1 if status_flag(combined, "improdutiva", "improdutivo", "improdutividade") else 0
    perda = 1 if any([cancelada, frustrada, reagendada, improdutiva]) else 0
    return concluida, cancelada, frustrada, reagendada, improdutiva, perda


def import_operation_workbook(path, mode="replace", df=None, sheet_name=None):
    init_db()
    path = Path(path)
    if df is None:
        df, best = read_best_excel_sheet(path)
        sheet_name = best["sheet"]
    df["_data"] = pd.to_datetime(df["data"], errors="coerce")
    df["_mes"] = df["_data"].dt.to_period("M").astype(str).replace("NaT", "Sem data")

    imported_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    rows = []
    for idx, source in df.iterrows():
        status = clean(source.get("status"), "Sem status")
        reason = clean(source.get("motivo_nao_realizada"), "Sem motivo")
        concluida, cancelada, frustrada, reagendada, improdutiva, perda = operation_flags(status, reason)
        record = {
            "row_index": idx,
            "date": clean_date(source["_data"]),
            "month": clean(source["_mes"], "Sem data"),
            "numero_assistencia": clean(source.get("numero_da_assistencia"), ""),
            "filial": clean(source.get("filial"), "Sem filial"),
            "categoria": clean(source.get("categoria"), "Sem categoria"),
            "cidade": clean(source.get("cidade"), "Sem cidade"),
            "estado": clean(source.get("estado"), "Sem UF"),
            "regional": clean(source.get("celula"), "Sem regional"),
            "area_trabalho": clean(source.get("area_de_trabalho"), "Sem area"),
            "hub": clean(source.get("hub"), "Sem HUB"),
            "posto": clean(source.get("posto"), "Sem posto"),
            "tipo_atividade": clean(source.get("tipo_atividade"), "Sem atividade"),
            "status": status,
            "recurso": clean(source.get("recurso"), "Sem recurso"),
            "nome_cliente": clean(source.get("nome_cliente"), "Sem nome"),
            "cliente": clean(source.get("cliente"), "Sem cliente"),
            "motivo": reason,
            "id_montador": clean(source.get("id_montador"), "Sem ID"),
            "concluida": concluida,
            "cancelada": cancelada,
            "frustrada": frustrada,
            "reagendada": reagendada,
            "improdutiva": improdutiva,
            "perda": perda,
            "source_file": path.name,
            "imported_at": imported_at,
        }
        record["record_hash"] = operation_row_hash(record)
        rows.append(record)

    with connect() as conn:
        if mode == "replace":
            conn.execute("DELETE FROM operations")
            conn.execute("DELETE FROM operation_batches")

        before = conn.execute("SELECT COUNT(*) AS total FROM operations").fetchone()["total"]
        conn.executemany(
            """
            INSERT INTO operations (
                record_hash, data, mes, numero_assistencia, filial, categoria, cidade, estado, regional,
                area_trabalho, hub, posto, tipo_atividade, status, recurso, nome_cliente, cliente,
                motivo_nao_realizada, id_montador, concluida, cancelada, frustrada, reagendada,
                improdutiva, perda, source_file, imported_at
            ) VALUES (
                :record_hash, :date, :month, :numero_assistencia, :filial, :categoria, :cidade, :estado, :regional,
                :area_trabalho, :hub, :posto, :tipo_atividade, :status, :recurso, :nome_cliente, :cliente,
                :motivo, :id_montador, :concluida, :cancelada, :frustrada, :reagendada,
                :improdutiva, :perda, :source_file, :imported_at
            )
            """,
            rows,
        )
        after = conn.execute("SELECT COUNT(*) AS total FROM operations").fetchone()["total"]
        inserted = after - before
        dates = [item["date"] for item in rows if item["date"]]
        conn.execute(
            """
            INSERT INTO operation_batches (
                file_name, imported_at, mode, rows_total, rows_inserted, concluidas, date_min, date_max
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                path.name,
                imported_at,
                mode,
                len(rows),
                inserted,
                sum(item["concluida"] for item in rows),
                min(dates) if dates else None,
                max(dates) if dates else None,
            ),
        )
        sync_geocidades(conn, force=True)
    return {
        "type": "operation",
        "file": path.name,
        "sheet": sheet_name,
        "mode": mode,
        "rowsTotal": len(rows),
        "rowsInserted": inserted,
        "concluidas": sum(item["concluida"] for item in rows),
        "dateMin": min(dates) if dates else None,
        "dateMax": max(dates) if dates else None,
    }


def import_workbook(path, mode="replace"):
    init_db()
    path = Path(path)
    df, best = read_best_excel_sheet(path)
    if best["type"] == "operation":
        return import_operation_workbook(path, mode=mode, df=df, sheet_name=best["sheet"])

    try:
        proc_col = find_col(df.columns, "tipo", "reclamacao", "csr")
        consolidated_col = find_col(df.columns, "procedencia", "consolidado")
    except KeyError as exc:
        columns = ", ".join(str(col) for col in df.columns)
        raise ValueError(
            f"A aba '{best['sheet']}' parece ser de reclamacoes, mas faltam campos obrigatorios. "
            f"Campo esperado: {exc}. Colunas encontradas: {columns}"
        ) from exc

    status_norm = df[proc_col].astype("string").str.strip().str.upper()
    consolidated = pd.to_numeric(df[consolidated_col], errors="coerce").fillna(0)
    df["_procedente"] = (status_norm.eq("PROCEDENTE") | consolidated.eq(1)).fillna(False)
    df["_data"] = pd.to_datetime(df["data_reclamacao"], errors="coerce")
    df["_dt_fechamento"] = pd.to_datetime(df["dt_fechamento"], errors="coerce")
    df["_mes"] = df["_data"].dt.to_period("M").astype(str).replace("NaT", "Sem data")
    df["_mes_fechamento"] = df["_dt_fechamento"].dt.to_period("M").astype(str).replace("NaT", "Sem data")

    imported_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    rows = []
    for idx, source in df.iterrows():
        record = {
            "row_index": idx,
            "id_feedback": clean(source.get("id_feedback"), ""),
            "id_assistservico": clean(source.get("id_assistservico"), ""),
            "id_assistencia": clean(source.get("id_assistencia"), ""),
            "date": clean_date(source["_data"]),
            "close_date": clean_date(source["_dt_fechamento"]),
            "month": clean(source["_mes"], "Sem data"),
            "close_month": clean(source["_mes_fechamento"], "Sem data"),
            "cliente": clean(source.get("Cliente"), "Sem cliente"),
            "cidade": clean(source.get("desc_cidade"), "Sem cidade"),
            "uf": clean(source.get("sigla"), "Sem UF"),
            "regiao": clean(source.get("regiao"), "Sem região"),
            "hub": clean(source.get("hub"), "Sem HUB"),
            "servico": clean(source.get("desc_servico"), "Sem serviço"),
            "status": clean(source.get("Statusreclamacao"), "Sem status"),
            "validacao": clean(source.get(proc_col), "Sem validação"),
            "motivo": clean(source.get("Motivo_Feedback"), "Sem motivo"),
            "grupo": clean(source.get("GrupoMotivoFeedback"), "Sem grupo"),
            "prestador": clean(source.get("Prestador"), "Sem prestador"),
            "analista": clean(source.get("analistasac"), "Sem analista"),
            "business_days_close": business_days_between(source["_data"], source["_dt_fechamento"]),
            "procedente": 1 if bool(source["_procedente"] is True or source["_procedente"] == 1) else 0,
            "source_file": path.name,
            "imported_at": imported_at,
        }
        record["sla_ok"] = (
            1
            if record["business_days_close"] is not None
            and record["business_days_close"] <= SLA_BUSINESS_DAYS_TARGET
            else 0
        )
        record["record_hash"] = row_hash(record)
        rows.append(record)

    with connect() as conn:
        if mode == "replace":
            conn.execute("DELETE FROM complaints")
            conn.execute("DELETE FROM import_batches")

        before = conn.execute("SELECT COUNT(*) AS total FROM complaints").fetchone()["total"]
        conn.executemany(
            """
            INSERT INTO complaints (
                record_hash, data_reclamacao, dt_fechamento, mes, mes_fechamento, cliente, cidade, uf, regiao, hub, servico,
                status_reclamacao, validacao, motivo, grupo_motivo, prestador, analista_sac, dias_uteis_fechamento, sla_ok, procedente,
                source_file, imported_at
            ) VALUES (
                :record_hash, :date, :close_date, :month, :close_month, :cliente, :cidade, :uf, :regiao, :hub, :servico,
                :status, :validacao, :motivo, :grupo, :prestador, :analista, :business_days_close, :sla_ok, :procedente,
                :source_file, :imported_at
            )
            """,
            rows,
        )
        after = conn.execute("SELECT COUNT(*) AS total FROM complaints").fetchone()["total"]
        inserted = after - before
        dates = [item["date"] for item in rows if item["date"]]
        conn.execute(
            """
            INSERT INTO import_batches (
                file_name, imported_at, mode, rows_total, rows_inserted, procedentes, date_min, date_max
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                path.name,
                imported_at,
                mode,
                len(rows),
                inserted,
                sum(item["procedente"] for item in rows),
                min(dates) if dates else None,
                max(dates) if dates else None,
            ),
        )
        sync_geocidades(conn, force=True)
    return {
        "type": "complaints",
        "file": path.name,
        "sheet": best["sheet"],
        "mode": mode,
        "rowsTotal": len(rows),
        "rowsInserted": inserted,
        "procedentes": sum(item["procedente"] for item in rows),
        "dateMin": min(dates) if dates else None,
        "dateMax": max(dates) if dates else None,
    }


def build_where(params, procedentes_only=False):
    where = []
    args = []
    if procedentes_only:
        where.append("procedente = 1")
    if params.get("regiao"):
        where.append("regiao = ?")
        args.append(params["regiao"])
    if params.get("hub"):
        where.append("hub = ?")
        args.append(params["hub"])
    if params.get("cliente"):
        where.append("cliente = ?")
        args.append(params["cliente"])
    if params.get("cidade"):
        where.append("cidade LIKE ?")
        args.append(f"%{params['cidade']}%")
    if params.get("analista"):
        where.append("analista_sac = ?")
        args.append(params["analista"])

    add_period_filter(where, args, "data_reclamacao", "mes", params, "_max_date")

    clause = " WHERE " + " AND ".join(where) if where else ""
    return clause, args


def build_close_where(params, closed_only=True):
    where = []
    args = []
    if closed_only:
        where.append("UPPER(status_reclamacao) = 'FECHADA'")
    where.append("dt_fechamento IS NOT NULL")
    if params.get("regiao"):
        where.append("regiao = ?")
        args.append(params["regiao"])
    if params.get("hub"):
        where.append("hub = ?")
        args.append(params["hub"])
    if params.get("cliente"):
        where.append("cliente = ?")
        args.append(params["cliente"])
    if params.get("cidade"):
        where.append("cidade LIKE ?")
        args.append(f"%{params['cidade']}%")
    if params.get("analista"):
        where.append("analista_sac = ?")
        args.append(params["analista"])

    if "_max_close_date" not in params and "_max_date" in params:
        params["_max_close_date"] = params["_max_date"]
    add_period_filter(where, args, "dt_fechamento", "mes_fechamento", params, "_max_close_date")

    clause = " WHERE " + " AND ".join(where) if where else ""
    return clause, args


def count_by(conn, params, field, limit=12):
    where, args = build_where(params, procedentes_only=params.get("scope", "procedentes") == "procedentes")
    rows = conn.execute(
        f"""
        SELECT {field} AS label, COUNT(*) AS value
        FROM complaints
        {where}
        GROUP BY {field}
        ORDER BY value DESC, label ASC
        LIMIT ?
        """,
        [*args, limit],
    ).fetchall()
    return [dict(row) for row in rows]


def table_combo(conn, params):
    where, args = build_where(params, procedentes_only=True)
    rows = conn.execute(
        f"""
        SELECT regiao, hub, cliente, cidade, COUNT(*) AS procedentes
        FROM complaints
        {where}
        GROUP BY regiao, hub, cliente, cidade
        ORDER BY procedentes DESC, regiao, hub, cliente, cidade
        LIMIT 50
        """,
        args,
    ).fetchall()
    return [dict(row) for row in rows]


def table_group_region(conn, params):
    where, args = build_where(params, procedentes_only=True)
    rows = conn.execute(
        f"""
        SELECT regiao, grupo_motivo AS grupo, COUNT(*) AS procedentes
        FROM complaints
        {where}
        GROUP BY regiao, grupo_motivo
        ORDER BY procedentes DESC, regiao, grupo
        LIMIT 40
        """,
        args,
    ).fetchall()
    return [dict(row) for row in rows]


def table_providers(conn, params):
    where, args = build_where(params, procedentes_only=True)
    rows = conn.execute(
        f"""
        SELECT prestador, hub, cidade, COUNT(*) AS procedentes
        FROM complaints
        {where}
        GROUP BY prestador, hub, cidade
        ORDER BY procedentes DESC, prestador
        LIMIT 40
        """,
        args,
    ).fetchall()
    return [dict(row) for row in rows]


def api_productivity(query):
    init_db()
    params = {key: values[0] for key, values in query.items() if values and values[0]}
    closed_weight = float(EFFICIENCY_SCORE_CONFIG.get("closed_weight", 1.0))
    sla_weight = float(EFFICIENCY_SCORE_CONFIG.get("sla_weight", 1.0))
    backlog_penalty = float(EFFICIENCY_SCORE_CONFIG.get("backlog_penalty", 0.75))
    with connect() as conn:
        max_date = conn.execute("SELECT MAX(data_reclamacao) AS value FROM complaints").fetchone()["value"]
        max_close_date = conn.execute("SELECT MAX(dt_fechamento) AS value FROM complaints").fetchone()["value"]
        params["_max_date"] = max_date
        params["_max_close_date"] = max_close_date
        where, args = build_where(params, procedentes_only=False)
        close_where, close_args = build_close_where(params, closed_only=True)

        total = conn.execute(f"SELECT COUNT(*) AS value FROM complaints {where}", args).fetchone()["value"]
        analysts = conn.execute(
            f"SELECT COUNT(DISTINCT analista_sac) AS value FROM complaints {where}", args
        ).fetchone()["value"]
        days = conn.execute(
            f"SELECT COUNT(DISTINCT data_reclamacao) AS value FROM complaints {where}", args
        ).fetchone()["value"]
        months = conn.execute(
            f"SELECT COUNT(DISTINCT CASE WHEN mes != 'Sem data' THEN mes END) AS value FROM complaints {where}",
            args,
        ).fetchone()["value"]
        closed_total = conn.execute(
            f"SELECT COUNT(*) AS value FROM complaints {close_where}", close_args
        ).fetchone()["value"]
        closed_days = conn.execute(
            f"SELECT COUNT(DISTINCT dt_fechamento) AS value FROM complaints {close_where}", close_args
        ).fetchone()["value"]
        closed_analysts = conn.execute(
            f"SELECT COUNT(DISTINCT analista_sac) AS value FROM complaints {close_where}", close_args
        ).fetchone()["value"]
        goal_expected = closed_analysts * closed_days * DAILY_CLOSED_GOAL
        sla = conn.execute(
            f"""
            SELECT
              COUNT(*) AS fechadas_com_prazo,
              SUM(CASE WHEN sla_ok = 1 THEN 1 ELSE 0 END) AS dentro_sla,
              ROUND(AVG(dias_uteis_fechamento), 2) AS media_dias_uteis
            FROM complaints
            {close_where}
            """,
            close_args,
        ).fetchone()
        backlog_rows = conn.execute(
            f"""
            SELECT status_reclamacao AS label, COUNT(*) AS value
            FROM complaints
            {where + (" AND UPPER(status_reclamacao) != 'FECHADA'" if where else " WHERE UPPER(status_reclamacao) != 'FECHADA'")}
            GROUP BY status_reclamacao
            ORDER BY value DESC, status_reclamacao ASC
            """
        , args).fetchall()
        status = conn.execute(
            f"""
            SELECT
              SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
              SUM(CASE WHEN UPPER(validacao) = 'IMPROCEDENTE' THEN 1 ELSE 0 END) AS improcedentes,
              SUM(CASE WHEN validacao = 'Sem validacao' OR validacao = 'Sem validação' OR validacao IS NULL THEN 1 ELSE 0 END) AS sem_status,
              SUM(CASE WHEN UPPER(status_reclamacao) = 'FECHADA' THEN 1 ELSE 0 END) AS fechadas
            FROM complaints
            {where}
            """,
            args,
        ).fetchone()
        lead = conn.execute(
            f"""
            SELECT analista_sac AS label, COUNT(*) AS value
            FROM complaints
            {where}
            GROUP BY analista_sac
            ORDER BY value DESC, analista_sac ASC
            LIMIT 1
            """,
            args,
        ).fetchone()

        analyst_rows = conn.execute(
            f"""
            SELECT
              analista_sac AS analista,
              COUNT(*) AS total,
              SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
              SUM(CASE WHEN UPPER(validacao) = 'IMPROCEDENTE' THEN 1 ELSE 0 END) AS improcedentes,
              SUM(CASE WHEN validacao = 'Sem validacao' OR validacao = 'Sem validação' OR validacao IS NULL THEN 1 ELSE 0 END) AS sem_status,
              SUM(CASE WHEN UPPER(status_reclamacao) = 'FECHADA' THEN 1 ELSE 0 END) AS fechadas,
              COUNT(DISTINCT data_reclamacao) AS dias_ativos,
              COUNT(DISTINCT CASE WHEN mes != 'Sem data' THEN mes END) AS meses_ativos,
              ROUND(COUNT(*) * 1.0 / NULLIF(COUNT(DISTINCT data_reclamacao), 0), 2) AS media_dia,
              ROUND(COUNT(*) * 1.0 / NULLIF(COUNT(DISTINCT CASE WHEN mes != 'Sem data' THEN mes END), 0), 2) AS media_mes,
              ROUND(SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) * 100.0 / COUNT(*), 1) AS taxa_procedencia
            FROM complaints
            {where}
            GROUP BY analista_sac
            ORDER BY total DESC, analista_sac ASC
            LIMIT 80
            """,
            args,
        ).fetchall()

        monthly_where = close_where + " AND mes_fechamento != 'Sem data'"
        monthly = conn.execute(
            f"""
            SELECT mes_fechamento AS label, COUNT(*) AS value
            FROM complaints
            {monthly_where}
            GROUP BY mes_fechamento
            ORDER BY mes_fechamento ASC
            """,
            close_args,
        ).fetchall()
        daily = conn.execute(
            f"""
            SELECT dt_fechamento AS label, COUNT(*) AS value
            FROM complaints
            {close_where}
            GROUP BY dt_fechamento
            ORDER BY dt_fechamento ASC
            LIMIT 120
            """,
            close_args,
        ).fetchall()
        status_by_analyst = conn.execute(
            f"""
            SELECT
              analista_sac AS analista,
              SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
              SUM(CASE WHEN UPPER(validacao) = 'IMPROCEDENTE' THEN 1 ELSE 0 END) AS improcedentes,
              SUM(CASE WHEN validacao = 'Sem validacao' OR validacao = 'Sem validação' OR validacao IS NULL THEN 1 ELSE 0 END) AS sem_status,
              SUM(CASE WHEN UPPER(status_reclamacao) = 'FECHADA' THEN 1 ELSE 0 END) AS fechadas,
              COUNT(*) AS total
            FROM complaints
            {where}
            GROUP BY analista_sac
            ORDER BY total DESC, analista_sac ASC
            LIMIT 30
            """,
            args,
        ).fetchall()
        closed_by_analyst = conn.execute(
            f"""
            SELECT analista_sac AS label, COUNT(*) AS value
            FROM complaints
            {close_where}
            GROUP BY analista_sac
            ORDER BY value DESC, analista_sac ASC
            LIMIT 15
            """,
            close_args,
        ).fetchall()
        closed_daily = conn.execute(
            f"""
            SELECT dt_fechamento AS label, COUNT(*) AS value
            FROM complaints
            {close_where}
            GROUP BY dt_fechamento
            ORDER BY dt_fechamento ASC
            LIMIT 120
            """,
            close_args,
        ).fetchall()
        daily_by_analyst = conn.execute(
            f"""
            SELECT dt_fechamento AS data, analista_sac AS analista, COUNT(*) AS total
            FROM complaints
            {close_where}
            GROUP BY dt_fechamento, analista_sac
            ORDER BY dt_fechamento DESC, total DESC, analista_sac ASC
            LIMIT 120
            """,
            close_args,
        ).fetchall()
        goal_by_analyst = conn.execute(
            f"""
            SELECT
              analista_sac AS analista,
              COUNT(*) AS fechadas,
              COUNT(DISTINCT dt_fechamento) AS dias_fechamento,
              COUNT(DISTINCT dt_fechamento) * {DAILY_CLOSED_GOAL} AS meta,
              ROUND(COUNT(*) * 100.0 / NULLIF(COUNT(DISTINCT dt_fechamento) * {DAILY_CLOSED_GOAL}, 0), 1) AS atingimento_meta
            FROM complaints
            {close_where}
            GROUP BY analista_sac
            ORDER BY atingimento_meta DESC, fechadas DESC, analista_sac ASC
            LIMIT 40
            """,
            close_args,
        ).fetchall()
        sla_by_analyst = conn.execute(
            f"""
            SELECT
              analista_sac AS analista,
              COUNT(*) AS fechadas,
              SUM(CASE WHEN sla_ok = 1 THEN 1 ELSE 0 END) AS dentro_sla,
              ROUND(SUM(CASE WHEN sla_ok = 1 THEN 1 ELSE 0 END) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa_sla,
              ROUND(AVG(dias_uteis_fechamento), 2) AS media_dias_uteis
            FROM complaints
            {close_where}
            GROUP BY analista_sac
            ORDER BY taxa_sla DESC, fechadas DESC, analista_sac ASC
            LIMIT 40
            """,
            close_args,
        ).fetchall()
        quality_volume = conn.execute(
            f"""
            SELECT
              analista_sac AS analista,
              COUNT(*) AS total,
              SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
              ROUND(SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa_procedencia,
              SUM(CASE WHEN UPPER(status_reclamacao) = 'FECHADA' THEN 1 ELSE 0 END) AS fechadas,
              ROUND(SUM(CASE WHEN sla_ok = 1 THEN 1 ELSE 0 END) * 100.0 / NULLIF(SUM(CASE WHEN UPPER(status_reclamacao) = 'FECHADA' THEN 1 ELSE 0 END), 0), 1) AS taxa_sla
            FROM complaints
            {where}
            GROUP BY analista_sac
            HAVING total >= 10
            ORDER BY total DESC, taxa_procedencia DESC
            LIMIT 40
            """,
            args,
        ).fetchall()
        open_where = where + (" AND UPPER(status_reclamacao) != 'FECHADA'" if where else " WHERE UPPER(status_reclamacao) != 'FECHADA'")
        backlog_aging = conn.execute(
            f"""
            SELECT faixa AS label, COUNT(*) AS value
            FROM (
              SELECT CASE
                WHEN data_reclamacao IS NULL THEN 'Sem data'
                WHEN julianday(?) - julianday(data_reclamacao) <= 5 THEN '0 a 5 dias'
                WHEN julianday(?) - julianday(data_reclamacao) <= 10 THEN '6 a 10 dias'
                WHEN julianday(?) - julianday(data_reclamacao) <= 20 THEN '11 a 20 dias'
                ELSE 'Mais de 20 dias'
              END AS faixa
              FROM complaints
              {open_where}
            )
            GROUP BY faixa
            ORDER BY CASE faixa
              WHEN '0 a 5 dias' THEN 1
              WHEN '6 a 10 dias' THEN 2
              WHEN '11 a 20 dias' THEN 3
              WHEN 'Mais de 20 dias' THEN 4
              ELSE 5
            END
            """,
            [max_date, max_date, max_date, *args],
        ).fetchall()
        backlog_by_analyst = conn.execute(
            f"""
            SELECT
              analista_sac AS analista,
              COUNT(*) AS backlog,
              SUM(CASE WHEN data_reclamacao IS NOT NULL AND julianday(?) - julianday(data_reclamacao) > {SLA_BUSINESS_DAYS_TARGET} THEN 1 ELSE 0 END) AS vencidas,
              SUM(CASE WHEN data_reclamacao IS NULL OR julianday(?) - julianday(data_reclamacao) <= {SLA_BUSINESS_DAYS_TARGET} THEN 1 ELSE 0 END) AS dentro_prazo
            FROM complaints
            {open_where}
            GROUP BY analista_sac
            ORDER BY backlog DESC, analista_sac ASC
            LIMIT 60
            """,
            [max_date, max_date, *args],
        ).fetchall()
        efficiency = conn.execute(
            f"""
            WITH base AS (
              SELECT
                analista_sac AS analista,
                COUNT(*) AS total,
                SUM(CASE WHEN UPPER(status_reclamacao) = 'FECHADA' THEN 1 ELSE 0 END) AS fechadas,
                SUM(CASE WHEN UPPER(status_reclamacao) != 'FECHADA' THEN 1 ELSE 0 END) AS backlog,
                SUM(CASE WHEN sla_ok = 1 THEN 1 ELSE 0 END) AS dentro_sla,
                SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
                ROUND(AVG(CASE WHEN dias_uteis_fechamento IS NOT NULL THEN dias_uteis_fechamento END), 2) AS media_dias_uteis
              FROM complaints
              {where}
              GROUP BY analista_sac
            )
            SELECT
              analista,
              total,
              fechadas,
              backlog,
              dentro_sla,
              ROUND(dentro_sla * 100.0 / NULLIF(fechadas, 0), 1) AS taxa_sla,
              ROUND(procedentes * 100.0 / NULLIF(total, 0), 1) AS taxa_procedencia,
              media_dias_uteis,
              ROUND(
                fechadas * {closed_weight}
                + COALESCE(dentro_sla * 100.0 / NULLIF(fechadas, 0), 0) * {sla_weight}
                - backlog * {backlog_penalty},
                1
              ) AS score_eficiencia
            FROM base
            WHERE total > 0
            ORDER BY score_eficiencia DESC, fechadas DESC, analista ASC
            LIMIT 80
            """,
            args,
        ).fetchall()
        sla_by_client = conn.execute(
            f"""
            SELECT
              cliente,
              COUNT(*) AS fechadas,
              SUM(CASE WHEN sla_ok = 1 THEN 1 ELSE 0 END) AS dentro_sla,
              ROUND(SUM(CASE WHEN sla_ok = 1 THEN 1 ELSE 0 END) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa_sla,
              ROUND(AVG(dias_uteis_fechamento), 2) AS media_dias_uteis
            FROM complaints
            {close_where}
            GROUP BY cliente
            HAVING fechadas >= 5
            ORDER BY taxa_sla ASC, fechadas DESC, cliente ASC
            LIMIT 50
            """,
            close_args,
        ).fetchall()
        sla_by_reason = conn.execute(
            f"""
            SELECT
              motivo,
              COUNT(*) AS fechadas,
              SUM(CASE WHEN sla_ok = 1 THEN 1 ELSE 0 END) AS dentro_sla,
              ROUND(SUM(CASE WHEN sla_ok = 1 THEN 1 ELSE 0 END) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa_sla,
              ROUND(AVG(dias_uteis_fechamento), 2) AS media_dias_uteis
            FROM complaints
            {close_where}
            GROUP BY motivo
            HAVING fechadas >= 5
            ORDER BY media_dias_uteis DESC, fechadas DESC, motivo ASC
            LIMIT 50
            """,
            close_args,
        ).fetchall()
        concentration_client = conn.execute(
            f"""
            SELECT analista_sac AS analista, cliente AS dimensao, COUNT(*) AS total
            FROM complaints
            {where}
            GROUP BY analista_sac, cliente
            ORDER BY total DESC, analista_sac ASC
            LIMIT 30
            """,
            args,
        ).fetchall()
        concentration_hub = conn.execute(
            f"""
            SELECT analista_sac AS analista, hub AS dimensao, COUNT(*) AS total
            FROM complaints
            {where}
            GROUP BY analista_sac, hub
            ORDER BY total DESC, analista_sac ASC
            LIMIT 30
            """,
            args,
        ).fetchall()
        concentration_reason = conn.execute(
            f"""
            SELECT analista_sac AS analista, motivo AS dimensao, COUNT(*) AS total
            FROM complaints
            {where}
            GROUP BY analista_sac, motivo
            ORDER BY total DESC, analista_sac ASC
            LIMIT 30
            """,
            args,
        ).fetchall()
        detail_rows = conn.execute(
            f"""
            SELECT
              data_reclamacao, dt_fechamento, cliente, cidade, uf, regiao, hub,
              status_reclamacao, validacao, motivo, analista_sac,
              dias_uteis_fechamento, sla_ok
            FROM complaints
            {where}
            ORDER BY COALESCE(dt_fechamento, data_reclamacao) DESC, analista_sac ASC
            LIMIT 200
            """,
            args,
        ).fetchall()
        compare_params = dict(params)
        for key in ["period", "month", "start", "end"]:
            compare_params.pop(key, None)
        compare_params["_max_close_date"] = max_close_date
        compare_where, compare_args = build_close_where(compare_params, closed_only=True)
        comparison = conn.execute(
            f"""
            SELECT
              SUM(CASE WHEN mes_fechamento = strftime('%Y-%m', date(?, 'start of month')) THEN 1 ELSE 0 END) AS mes_atual,
              SUM(CASE WHEN mes_fechamento = strftime('%Y-%m', date(?, 'start of month', '-1 month')) THEN 1 ELSE 0 END) AS mes_anterior
            FROM complaints
            {compare_where}
            """,
            [max_close_date, max_close_date, *compare_args],
        ).fetchone()
        top_backlog = dict(backlog_rows[0]) if backlog_rows else None
        executive_summary = [
            f"Total fechado no recorte: {closed_total} reclamacoes.",
            f"Meta diaria: {DAILY_CLOSED_GOAL} fechadas por analista; atingimento geral de {round((closed_total / goal_expected * 100), 1) if goal_expected else 0}%.",
            f"SLA de ate {SLA_BUSINESS_DAYS_TARGET} dias uteis: {round(((sla['dentro_sla'] or 0) / (sla['fechadas_com_prazo'] or 1)) * 100, 1) if sla else 0}% dentro do prazo.",
            f"Backlog principal: {top_backlog['label']} com {top_backlog['value']} reclamacoes." if top_backlog else "Sem backlog no recorte.",
        ]

        return {
            "kpis": {
                "total": total,
                "analysts": analysts,
                "days": days,
                "months": months,
                "avgPerAnalyst": total / analysts if analysts else 0,
                "avgPerDay": total / days if days else 0,
                "avgPerMonth": total / months if months else 0,
                "procedentes": status["procedentes"] or 0,
                "improcedentes": status["improcedentes"] or 0,
                "semStatus": status["sem_status"] or 0,
                "fechadas": closed_total,
                "dailyGoal": DAILY_CLOSED_GOAL,
                "goalExpected": goal_expected,
                "goalPct": closed_total / goal_expected if goal_expected else 0,
                "slaTarget": SLA_BUSINESS_DAYS_TARGET,
                "slaOk": sla["dentro_sla"] or 0,
                "slaTotal": sla["fechadas_com_prazo"] or 0,
                "slaPct": (sla["dentro_sla"] or 0) / (sla["fechadas_com_prazo"] or 1),
                "avgBusinessDaysToClose": sla["media_dias_uteis"] or 0,
                "backlog": sum(row["value"] for row in backlog_rows),
                "comparison": {
                    "currentMonth": comparison["mes_atual"] or 0,
                    "previousMonth": comparison["mes_anterior"] or 0,
                    "delta": (comparison["mes_atual"] or 0) - (comparison["mes_anterior"] or 0),
                    "deltaPct": ((comparison["mes_atual"] or 0) - (comparison["mes_anterior"] or 0)) / (comparison["mes_anterior"] or 1),
                },
                "leadAnalyst": dict(lead) if lead else None,
            },
            "executiveSummary": executive_summary,
            "scoreConfig": {
                "closedWeight": closed_weight,
                "slaWeight": sla_weight,
                "backlogPenalty": backlog_penalty,
                "formula": (
                    f"Score eficiencia = fechadas x {closed_weight:g} + % SLA x {sla_weight:g} "
                    f"- backlog x {backlog_penalty:g}."
                ),
            },
            "charts": {
                "analystVolume": [{"label": row["analista"], "value": row["total"]} for row in analyst_rows[:15]],
                "monthly": [dict(row) for row in monthly],
                "daily": [dict(row) for row in daily],
                "closedDaily": [dict(row) for row in closed_daily],
                "closedByAnalyst": [dict(row) for row in closed_by_analyst],
                "backlog": [dict(row) for row in backlog_rows],
                "backlogAging": [dict(row) for row in backlog_aging],
                "statusByAnalyst": [dict(row) for row in status_by_analyst],
            },
            "tables": {
                "analysts": [dict(row) for row in analyst_rows],
                "dailyByAnalyst": [dict(row) for row in daily_by_analyst],
                "goalByAnalyst": [dict(row) for row in goal_by_analyst],
                "slaByAnalyst": [dict(row) for row in sla_by_analyst],
                "qualityVolume": [dict(row) for row in quality_volume],
                "efficiency": [dict(row) for row in efficiency],
                "backlogByAnalyst": [dict(row) for row in backlog_by_analyst],
                "slaByClient": [dict(row) for row in sla_by_client],
                "slaByReason": [dict(row) for row in sla_by_reason],
                "concentrationClient": [dict(row) for row in concentration_client],
                "concentrationHub": [dict(row) for row in concentration_hub],
                "concentrationReason": [dict(row) for row in concentration_reason],
                "details": [dict(row) for row in detail_rows],
            },
        }


def api_summary(query):
    init_db()
    params = {key: values[0] for key, values in query.items() if values and values[0]}
    with connect() as conn:
        max_date = conn.execute("SELECT MAX(data_reclamacao) AS value FROM complaints").fetchone()["value"]
        params["_max_date"] = max_date
        all_where, all_args = build_where(params, procedentes_only=False)
        proc_where, proc_args = build_where(params, procedentes_only=True)

        total = conn.execute(f"SELECT COUNT(*) AS value FROM complaints {all_where}", all_args).fetchone()["value"]
        procedentes = conn.execute(f"SELECT COUNT(*) AS value FROM complaints {proc_where}", proc_args).fetchone()["value"]
        cities = conn.execute(f"SELECT COUNT(DISTINCT cidade) AS value FROM complaints {proc_where}", proc_args).fetchone()["value"]
        hubs = conn.execute(f"SELECT COUNT(DISTINCT hub) AS value FROM complaints {proc_where}", proc_args).fetchone()["value"]
        lead = conn.execute(
            f"""
            SELECT cliente AS label, COUNT(*) AS value
            FROM complaints {proc_where}
            GROUP BY cliente
            ORDER BY value DESC, cliente ASC
            LIMIT 1
            """,
            proc_args,
        ).fetchone()
        monthly_where = proc_where + (" AND mes != 'Sem data'" if proc_where else " WHERE mes != 'Sem data'")
        monthly = conn.execute(
            f"""
            SELECT mes AS label, COUNT(*) AS value
            FROM complaints
            {monthly_where}
            GROUP BY mes
            ORDER BY mes ASC
            """,
            proc_args,
        ).fetchall()
        monthly_mix = conn.execute(
            f"""
            SELECT
              mes AS label,
              COUNT(*) AS total,
              SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
              SUM(CASE WHEN UPPER(validacao) = 'IMPROCEDENTE' THEN 1 ELSE 0 END) AS improcedentes,
              SUM(CASE WHEN validacao = 'Sem validacao' OR validacao = 'Sem validação' OR validacao IS NULL THEN 1 ELSE 0 END) AS sem_status,
              ROUND(SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa_procedencia
            FROM complaints
            {all_where + (" AND mes != 'Sem data'" if all_where else " WHERE mes != 'Sem data'")}
            GROUP BY mes
            ORDER BY mes ASC
            """,
            all_args,
        ).fetchall()
        validation_mix = conn.execute(
            f"""
            SELECT validacao AS label, COUNT(*) AS value
            FROM complaints
            {all_where}
            GROUP BY validacao
            ORDER BY value DESC, validacao ASC
            """
        , all_args).fetchall()
        status_mix = conn.execute(
            f"""
            SELECT status_reclamacao AS label, COUNT(*) AS value
            FROM complaints
            {all_where}
            GROUP BY status_reclamacao
            ORDER BY value DESC, status_reclamacao ASC
            LIMIT 10
            """
        , all_args).fetchall()
        rate_region = conn.execute(
            f"""
            SELECT regiao AS label, COUNT(*) AS total,
              SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
              ROUND(SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa
            FROM complaints
            {all_where}
            GROUP BY regiao
            HAVING total >= 10
            ORDER BY taxa DESC, total DESC
            LIMIT 15
            """
        , all_args).fetchall()
        rate_hub = conn.execute(
            f"""
            SELECT hub AS label, COUNT(*) AS total,
              SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
              ROUND(SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa
            FROM complaints
            {all_where}
            GROUP BY hub
            HAVING total >= 10
            ORDER BY taxa DESC, total DESC
            LIMIT 20
            """
        , all_args).fetchall()
        rate_client = conn.execute(
            f"""
            SELECT cliente AS label, COUNT(*) AS total,
              SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
              ROUND(SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa
            FROM complaints
            {all_where}
            GROUP BY cliente
            HAVING total >= 10
            ORDER BY taxa DESC, total DESC
            LIMIT 20
            """
        , all_args).fetchall()
        proc_details = conn.execute(
            f"""
            SELECT data_reclamacao, dt_fechamento, cliente, cidade, uf, regiao, hub, status_reclamacao,
                   validacao, motivo, grupo_motivo, analista_sac, dias_uteis_fechamento, sla_ok
            FROM complaints
            {proc_where}
            ORDER BY data_reclamacao DESC, cliente, hub
            LIMIT 200
            """
        , proc_args).fetchall()
        pareto_reason = conn.execute(
            f"""
            SELECT motivo AS label, COUNT(*) AS value
            FROM complaints
            {proc_where}
            GROUP BY motivo
            ORDER BY value DESC, motivo ASC
            LIMIT 20
            """
        , proc_args).fetchall()
        heatmap_region_reason = conn.execute(
            f"""
            SELECT regiao, motivo, COUNT(*) AS value
            FROM complaints
            {proc_where}
            GROUP BY regiao, motivo
            ORDER BY value DESC, regiao, motivo
            LIMIT 80
            """
        , proc_args).fetchall()
        recurrence_clients = conn.execute(
            f"""
            WITH ordered AS (
              SELECT cliente, data_reclamacao,
                     LAG(data_reclamacao) OVER (PARTITION BY cliente ORDER BY data_reclamacao) AS prev_date
              FROM complaints
              {all_where + (" AND data_reclamacao IS NOT NULL" if all_where else " WHERE data_reclamacao IS NOT NULL")}
            )
            SELECT cliente AS label, COUNT(*) AS reincidencias
            FROM ordered
            WHERE prev_date IS NOT NULL AND julianday(data_reclamacao) - julianday(prev_date) <= 30
            GROUP BY cliente
            ORDER BY reincidencias DESC, cliente ASC
            LIMIT 20
            """
        , all_args).fetchall()
        recurrence_total = len(recurrence_clients)
        provider_critical = conn.execute(
            f"""
            SELECT
              prestador,
              hub,
              cidade,
              COUNT(*) AS procedentes,
              SUM(CASE WHEN sla_ok = 0 AND dias_uteis_fechamento IS NOT NULL THEN 1 ELSE 0 END) AS sla_estourado,
              COUNT(*) + SUM(CASE WHEN sla_ok = 0 AND dias_uteis_fechamento IS NOT NULL THEN 1 ELSE 0 END) AS score
            FROM complaints
            {proc_where}
            GROUP BY prestador, hub, cidade
            HAVING procedentes > 0
            ORDER BY score DESC, procedentes DESC, prestador ASC
            LIMIT 50
            """
        , proc_args).fetchall()
        compare_params = dict(params)
        for key in ["period", "month", "start", "end"]:
            compare_params.pop(key, None)
        compare_params["_max_date"] = max_date
        compare_where, compare_args = build_where(compare_params, procedentes_only=False)
        comparison = conn.execute(
            f"""
            SELECT
              SUM(CASE WHEN mes = strftime('%Y-%m', date(?, 'start of month')) AND procedente = 1 THEN 1 ELSE 0 END) AS mes_atual,
              SUM(CASE WHEN mes = strftime('%Y-%m', date(?, 'start of month', '-1 month')) AND procedente = 1 THEN 1 ELSE 0 END) AS mes_anterior
            FROM complaints
            {compare_where}
            """,
            [max_date, max_date, *compare_args],
        ).fetchone()
        date_bounds = conn.execute(
            "SELECT MIN(data_reclamacao) AS minDate, MAX(data_reclamacao) AS maxDate FROM complaints"
        ).fetchone()
        batches = conn.execute(
            """
            SELECT file_name, imported_at, mode, rows_total, rows_inserted, procedentes, date_min, date_max
            FROM import_batches
            ORDER BY id DESC
            LIMIT 8
            """
        ).fetchall()

        return {
            "kpis": {
                "total": total,
                "procedentes": procedentes,
                "procedenciaPct": procedentes / total if total else 0,
                "improcedentes": sum(row["value"] for row in validation_mix if str(row["label"]).upper() == "IMPROCEDENTE"),
                "semStatus": sum(row["value"] for row in validation_mix if str(row["label"]).lower().startswith("sem")),
                "reincidencias": recurrence_total,
                "reincidenciaPct": recurrence_total / total if total else 0,
                "comparison": {
                    "currentMonth": comparison["mes_atual"] or 0,
                    "previousMonth": comparison["mes_anterior"] or 0,
                    "delta": (comparison["mes_atual"] or 0) - (comparison["mes_anterior"] or 0),
                    "deltaPct": ((comparison["mes_atual"] or 0) - (comparison["mes_anterior"] or 0)) / (comparison["mes_anterior"] or 1),
                },
                "cities": cities,
                "hubs": hubs,
                "leadCliente": dict(lead) if lead else None,
            },
            "dateBounds": dict(date_bounds),
            "charts": {
                "monthly": [dict(row) for row in monthly],
                "regiao": count_by(conn, params, "regiao", 8),
                "hub": count_by(conn, params, "hub", 12),
                "cliente": count_by(conn, params, "cliente", 10),
                "cidade": count_by(conn, params, "cidade", 15),
                "motivo": count_by(conn, params, "motivo", 15),
                "validationMix": [dict(row) for row in validation_mix],
                "statusMix": [dict(row) for row in status_mix],
                "paretoReason": [dict(row) for row in pareto_reason],
                "heatmapRegionReason": [dict(row) for row in heatmap_region_reason],
            },
            "tables": {
                "combo": table_combo(conn, params),
                "groupRegion": table_group_region(conn, params),
                "providers": table_providers(conn, params),
                "monthlyMix": [dict(row) for row in monthly_mix],
                "rateRegion": [dict(row) for row in rate_region],
                "rateHub": [dict(row) for row in rate_hub],
                "rateClient": [dict(row) for row in rate_client],
                "details": [dict(row) for row in proc_details],
                "recurrenceClients": [dict(row) for row in recurrence_clients],
                "providerCritical": [dict(row) for row in provider_critical],
            },
            "imports": [dict(row) for row in batches],
        }


def build_operation_where(params):
    where = []
    args = []
    if params.get("regiao"):
        where.append("regional = ?")
        args.append(params["regiao"])
    if params.get("hub"):
        where.append("hub = ?")
        args.append(params["hub"])
    if params.get("cliente"):
        where.append("cliente = ?")
        args.append(params["cliente"])
    if params.get("cidade"):
        where.append("cidade LIKE ?")
        args.append(f"%{params['cidade']}%")

    add_period_filter(where, args, "data", "mes", params, "_max_operation_date")

    clause = " WHERE " + " AND ".join(where) if where else ""
    return clause, args


def operation_count_by(conn, params, field, value_field="COUNT(*)", limit=12, order_by="value DESC"):
    where, args = build_operation_where(params)
    rows = conn.execute(
        f"""
        SELECT {field} AS label, {value_field} AS value
        FROM operations
        {where}
        GROUP BY {field}
        ORDER BY {order_by}, label ASC
        LIMIT ?
        """,
        [*args, limit],
    ).fetchall()
    return [dict(row) for row in rows]


def operation_performance_table(conn, params, field, limit=40):
    where, args = build_operation_where(params)
    rows = conn.execute(
        f"""
        SELECT
          {field} AS label,
          COUNT(*) AS total,
          SUM(concluida) AS concluidas,
          SUM(cancelada) AS canceladas,
          SUM(frustrada) AS frustradas,
          SUM(reagendada) AS reagendadas,
          SUM(improdutiva) AS improdutivas,
          SUM(perda) AS perdas,
          ROUND(SUM(concluida) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa_sucesso,
          ROUND(SUM(perda) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa_perda,
          ROUND((SUM(concluida) - SUM(perda)) * 100.0 / NULLIF(COUNT(*), 0), 1) AS score_operacional
        FROM operations
        {where}
        GROUP BY {field}
        HAVING total > 0
        ORDER BY taxa_perda DESC, total DESC, label ASC
        LIMIT ?
        """,
        [*args, limit],
    ).fetchall()
    return [dict(row) for row in rows]


def api_operation(query):
    init_db()
    params = {key: values[0] for key, values in query.items() if values and values[0]}
    with connect() as conn:
        max_date = conn.execute("SELECT MAX(data) AS value FROM operations").fetchone()["value"]
        params["_max_operation_date"] = max_date
        where, args = build_operation_where(params)
        totals = conn.execute(
            f"""
            SELECT
              COUNT(*) AS total,
              SUM(concluida) AS concluidas,
              SUM(cancelada) AS canceladas,
              SUM(frustrada) AS frustradas,
              SUM(reagendada) AS reagendadas,
              SUM(improdutiva) AS improdutivas,
              SUM(perda) AS perdas
            FROM operations
            {where}
            """,
            args,
        ).fetchone()
        total = totals["total"] or 0
        concluidas = totals["concluidas"] or 0
        perdas = totals["perdas"] or 0
        monthly = conn.execute(
            f"""
            SELECT mes AS label, COUNT(*) AS total, SUM(concluida) AS concluidas, SUM(perda) AS perdas,
                   ROUND(SUM(concluida) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa_sucesso
            FROM operations
            {where + (" AND mes != 'Sem data'" if where else " WHERE mes != 'Sem data'")}
            GROUP BY mes
            ORDER BY mes ASC
            """,
            args,
        ).fetchall()
        compare_params = dict(params)
        for key in ["period", "month", "start", "end"]:
            compare_params.pop(key, None)
        compare_params["_max_operation_date"] = max_date
        compare_where, compare_args = build_operation_where(compare_params)
        comparison = conn.execute(
            f"""
            SELECT
              SUM(CASE WHEN mes = strftime('%Y-%m', date(?, 'start of month')) THEN concluida ELSE 0 END) AS mes_atual,
              SUM(CASE WHEN mes = strftime('%Y-%m', date(?, 'start of month', '-1 month')) THEN concluida ELSE 0 END) AS mes_anterior
            FROM operations
            {compare_where}
            """,
            [max_date, max_date, *compare_args],
        ).fetchone()
        hub_table = operation_performance_table(conn, params, "hub", 40)
        regional_table = operation_performance_table(conn, params, "regional", 40)
        resource_table = operation_performance_table(conn, params, "recurso", 80)
        city_table = operation_performance_table(conn, params, "cidade", 60)
        service_table = operation_performance_table(conn, params, "tipo_atividade", 40)
        reason_rows = conn.execute(
            f"""
            SELECT motivo_nao_realizada AS motivo, COUNT(*) AS total,
                   SUM(CASE WHEN perda = 1 THEN 1 ELSE 0 END) AS perdas
            FROM operations
            {where + (" AND perda = 1" if where else " WHERE perda = 1")}
            GROUP BY motivo_nao_realizada
            ORDER BY perdas DESC, motivo ASC
            LIMIT 30
            """,
            args,
        ).fetchall()
        top_reason = reason_rows[0] if reason_rows else None
        worst_hub = hub_table[0] if hub_table else None
        critical_regional = regional_table[0] if regional_table else None
        critical_resource = resource_table[0] if resource_table else None
        best_resources = sorted(resource_table, key=lambda row: (-row["taxa_sucesso"], -row["total"], row["label"]))[:20]
        capacity_rows = conn.execute(
            f"""
            SELECT recurso AS label, COUNT(*) AS atendimentos, COUNT(DISTINCT data) AS dias_ativos,
                   ROUND(COUNT(*) * 1.0 / NULLIF(COUNT(DISTINCT data), 0), 2) AS media_dia
            FROM operations
            {where}
            GROUP BY recurso
            ORDER BY atendimentos DESC, label ASC
            LIMIT 60
            """,
            args,
        ).fetchall()
        imports = conn.execute(
            """
            SELECT file_name, imported_at, mode, rows_total, rows_inserted, concluidas, date_min, date_max
            FROM operation_batches
            ORDER BY id DESC
            LIMIT 8
            """
        ).fetchall()
        date_bounds = conn.execute(
            "SELECT MIN(data) AS minDate, MAX(data) AS maxDate FROM operations"
        ).fetchone()
        insights = [
            f"Taxa de sucesso operacional: {round(pct_value(concluidas, total) * 100, 1)}% ({concluidas} concluidas de {total} assistencias).",
            f"Taxa de perda operacional: {round(pct_value(perdas, total) * 100, 1)}% considerando canceladas, frustradas, reagendadas e improdutivas.",
        ]
        if worst_hub:
            insights.append(f"HUB com maior taxa de perda no recorte: {worst_hub['label']} ({worst_hub['taxa_perda']}%).")
        if top_reason:
            insights.append(f"Principal motivo de nao realizacao: {top_reason['motivo']} ({top_reason['perdas']} perdas).")
        return {
            "kpis": {
                "total": total,
                "concluidas": concluidas,
                "canceladas": totals["canceladas"] or 0,
                "frustradas": totals["frustradas"] or 0,
                "reagendadas": totals["reagendadas"] or 0,
                "improdutivas": totals["improdutivas"] or 0,
                "perdas": perdas,
                "successPct": pct_value(concluidas, total),
                "lossPct": pct_value(perdas, total),
                "worstHub": worst_hub,
                "criticalRegional": critical_regional,
                "criticalResource": critical_resource,
                "topReason": dict(top_reason) if top_reason else None,
                "comparison": {
                    "currentMonth": comparison["mes_atual"] or 0,
                    "previousMonth": comparison["mes_anterior"] or 0,
                    "delta": (comparison["mes_atual"] or 0) - (comparison["mes_anterior"] or 0),
                    "deltaPct": ((comparison["mes_atual"] or 0) - (comparison["mes_anterior"] or 0)) / (comparison["mes_anterior"] or 1),
                },
            },
            "dateBounds": dict(date_bounds),
            "insights": insights,
            "charts": {
                "monthly": [{"label": row["label"], "value": row["concluidas"] or 0} for row in monthly],
                "lossMonthly": [{"label": row["label"], "value": row["perdas"] or 0} for row in monthly],
                "status": operation_count_by(conn, params, "status", "COUNT(*)", 12),
                "hub": [{"label": row["label"], "value": row["taxa_perda"] or 0} for row in hub_table[:15]],
                "regional": [{"label": row["label"], "value": row["taxa_perda"] or 0} for row in regional_table[:12]],
                "resource": [{"label": row["label"], "value": row["taxa_perda"] or 0} for row in resource_table[:15]],
                "reason": [{"label": row["motivo"], "value": row["perdas"] or 0} for row in reason_rows],
                "cityLoss": [{"label": row["label"], "value": row["taxa_perda"] or 0} for row in city_table[:20]],
            },
            "tables": {
                "hubPerformance": hub_table,
                "regionalPerformance": regional_table,
                "resourceBest": best_resources,
                "resourceCritical": resource_table[:40],
                "reasonPareto": [dict(row) for row in reason_rows],
                "cityRisk": city_table,
                "servicePerformance": service_table,
                "capacity": [dict(row) for row in capacity_rows],
            },
            "imports": [dict(row) for row in imports],
        }


def risk_status(score):
    if score >= 70:
        return "Vermelho"
    if score >= 40:
        return "Amarelo"
    return "Verde"


def api_score_dimension(query, dimension):
    init_db()
    fields = {
        "prestadores": ("prestador", "prestador"),
        "hubs": ("hub", "hub"),
        "cidades": ("cidade || ' / ' || uf", "cidade"),
        "clientes": ("cliente", "cliente"),
    }
    if dimension not in fields:
        return {"rows": []}
    select_expr, label_field = fields[dimension]
    params = {key: values[0] for key, values in query.items() if values and values[0]}
    weights = CONFIG.get("risk_score", DEFAULT_CONFIG["risk_score"])
    volume_weight = float(weights.get("volume_weight", 0.40))
    procedencia_weight = float(weights.get("procedencia_weight", 0.30))
    sla_weight = float(weights.get("sla_weight", 0.20))
    reincidencia_weight = float(weights.get("reincidencia_weight", 0.10))
    with connect() as conn:
        max_date = conn.execute("SELECT MAX(data_reclamacao) AS value FROM complaints").fetchone()["value"]
        params["_max_date"] = max_date
        where, args = build_where(params, procedentes_only=False)
        rows = conn.execute(
            f"""
            WITH base AS (
              SELECT
                {select_expr} AS label,
                COUNT(*) AS reclamacoes,
                SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
                SUM(CASE WHEN UPPER(validacao) = 'IMPROCEDENTE' THEN 1 ELSE 0 END) AS improcedentes,
                SUM(CASE WHEN UPPER(status_reclamacao) = 'FECHADA' THEN 1 ELSE 0 END) AS fechadas,
                SUM(CASE WHEN UPPER(status_reclamacao) = 'FECHADA' AND sla_ok = 0 THEN 1 ELSE 0 END) AS fora_sla
              FROM complaints
              {where}
              GROUP BY label
            ),
            ranked AS (
              SELECT
                *,
                MAX(reclamacoes) OVER () AS max_reclamacoes
              FROM base
            )
            SELECT
              label,
              reclamacoes,
              procedentes,
              improcedentes,
              fechadas,
              fora_sla,
              ROUND(procedentes * 100.0 / NULLIF(reclamacoes, 0), 1) AS taxa_procedencia,
              ROUND(fora_sla * 100.0 / NULLIF(fechadas, 0), 1) AS taxa_fora_sla,
              ROUND(
                ((reclamacoes * 100.0 / NULLIF(max_reclamacoes, 0)) * {volume_weight})
                + (COALESCE(procedentes * 100.0 / NULLIF(reclamacoes, 0), 0) * {procedencia_weight})
                + (COALESCE(fora_sla * 100.0 / NULLIF(fechadas, 0), 0) * {sla_weight})
                + (0 * {reincidencia_weight}),
                1
              ) AS score_risco
            FROM ranked
            WHERE label IS NOT NULL
            ORDER BY score_risco DESC, reclamacoes DESC, label ASC
            LIMIT 100
            """,
            args,
        ).fetchall()
    output = []
    for row in rows:
        item = dict(row)
        item["status_risco"] = risk_status(item["score_risco"] or 0)
        item["dimensao"] = label_field
        item["formula"] = "40% volume + 30% taxa de procedencia + 20% SLA estourado + 10% reincidencia"
        output.append(item)
    return {"dimension": dimension, "rows": output}


def city_key(city, uf):
    return f"{text_key(city)}|{text_key(normalize_uf(uf))}"


def load_static_city_coordinates():
    if not CITIES_PATH.exists():
        return []
    try:
        rows = json.loads(CITIES_PATH.read_text(encoding="utf-8-sig"))
    except Exception:
        return []
    return rows if isinstance(rows, list) else []


def load_city_coordinates():
    rows = load_static_city_coordinates()
    coords = {}
    for row in rows:
        coords[city_key(row.get("cidade"), row.get("uf"))] = row
    return coords


def upsert_geocity(conn, cidade, uf, lat=None, lon=None, fonte="base", uf_corrigida_de=None):
    cidade = clean(cidade, "Sem cidade")
    cidade_normalizada = normalizar_cidade(cidade)
    uf = normalize_uf(uf)
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    cidade_key = text_key(cidade)
    uf_key = text_key(uf)
    has_coords = lat is not None and lon is not None
    conn.execute(
        """
        INSERT INTO geocidades (
            cidade, uf, cidade_key, uf_key, lat, lon, fonte, precisa_coordenada,
            uf_corrigida_de, created_at, updated_at
        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ON CONFLICT(cidade_key, uf_key) DO UPDATE SET
            cidade = excluded.cidade,
            uf = excluded.uf,
            lat = COALESCE(excluded.lat, geocidades.lat),
            lon = COALESCE(excluded.lon, geocidades.lon),
            fonte = CASE
                WHEN excluded.lat IS NOT NULL AND excluded.lon IS NOT NULL THEN excluded.fonte
                ELSE geocidades.fonte
            END,
            precisa_coordenada = CASE
                WHEN COALESCE(excluded.lat, geocidades.lat) IS NOT NULL
                 AND COALESCE(excluded.lon, geocidades.lon) IS NOT NULL THEN 0
                ELSE 1
            END,
            uf_corrigida_de = COALESCE(excluded.uf_corrigida_de, geocidades.uf_corrigida_de),
            updated_at = excluded.updated_at
        """,
        (
            cidade_normalizada or cidade,
            uf,
            cidade_key,
            uf_key,
            lat,
            lon,
            fonte,
            0 if has_coords else 1,
            uf_corrigida_de,
            now,
            now,
        ),
    )


def seed_geocidades_from_static(conn):
    for row in load_static_city_coordinates():
        upsert_geocity(
            conn,
            row.get("cidade"),
            row.get("uf"),
            row.get("lat"),
            row.get("lon"),
            fonte="static/data/cidades_brasil.json",
        )


def upsert_pending_geocity(conn, cidade, uf, uf_original, origem, reclamacoes=0, procedentes=0, assistencias=0):
    cidade_normalizada = normalizar_cidade(cidade)
    if not cidade_normalizada:
        return
    uf = normalize_uf(uf)
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    ocorrencias = (reclamacoes or 0) + (assistencias or 0)
    conn.execute(
        """
        INSERT INTO geocidades_pendentes (
            cidade, uf, uf_original, cidade_normalizada, ocorrencias, origem,
            reclamacoes, procedentes, assistencias, ignorada, created_at, updated_at
        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, 0, ?, ?)
        ON CONFLICT(cidade_normalizada, uf, origem) DO UPDATE SET
            cidade = excluded.cidade,
            uf_original = excluded.uf_original,
            ocorrencias = excluded.ocorrencias,
            reclamacoes = excluded.reclamacoes,
            procedentes = excluded.procedentes,
            assistencias = excluded.assistencias,
            updated_at = excluded.updated_at
        """,
        (
            cidade_normalizada,
            uf,
            uf_original,
            cidade_normalizada,
            ocorrencias,
            origem,
            reclamacoes or 0,
            procedentes or 0,
            assistencias or 0,
            now,
            now,
        ),
    )


def remove_pending_geocity(conn, cidade, uf):
    conn.execute(
        "DELETE FROM geocidades_pendentes WHERE cidade_normalizada = ? AND uf = ?",
        (normalizar_cidade(cidade), normalize_uf(uf)),
    )


def infer_city_uf(conn, cidade, uf):
    uf = normalize_uf(uf)
    if not is_unknown_uf(uf):
        return uf, None
    cidade_key = text_key(cidade)
    rows = conn.execute(
        """
        SELECT uf
        FROM geocidades
        WHERE cidade_key = ? AND uf_key != 'ni'
        GROUP BY uf
        """,
        (cidade_key,),
    ).fetchall()
    candidates = {normalize_uf(row["uf"]) for row in rows if not is_unknown_uf(row["uf"])}
    if len(candidates) == 1:
        return next(iter(candidates)), uf
    source_rows = conn.execute(
        """
        SELECT uf AS value FROM complaints WHERE cidade = ?
        UNION ALL
        SELECT estado AS value FROM operations WHERE cidade = ?
        """,
        (cidade, cidade),
    ).fetchall()
    source_candidates = {normalize_uf(row["value"]) for row in source_rows if not is_unknown_uf(row["value"])}
    if len(source_candidates) == 1:
        return next(iter(source_candidates)), uf
    return "NI", None


def sync_geocidades(conn, force=False):
    global GEO_SYNC_LAST_RUN
    with GEO_SYNC_LOCK:
        now_ts = time.time()
        if not force and GEO_SYNC_LAST_RUN and now_ts - GEO_SYNC_LAST_RUN < GEO_SYNC_TTL_SECONDS:
            return {"totalSeen": 0, "ufNiCorrected": 0, "skipped": True}
        seed_geocidades_from_static(conn)
        rows = conn.execute(
            """
            SELECT cidade, uf FROM complaints WHERE cidade IS NOT NULL AND cidade != 'Sem cidade'
            UNION
            SELECT cidade, estado AS uf FROM operations WHERE cidade IS NOT NULL AND cidade != 'Sem cidade'
            """
        ).fetchall()
        known_rows = conn.execute(
            """
            SELECT cidade_key, uf
            FROM geocidades
            WHERE uf_key != 'ni'
            GROUP BY cidade_key, uf
            """
        ).fetchall()
        known_by_city = {}
        for row in known_rows:
            known_by_city.setdefault(row["cidade_key"], set()).add(normalize_uf(row["uf"]))
        source_by_city = {}
        for row in rows:
            uf = normalize_uf(row["uf"])
            if not is_unknown_uf(uf):
                source_by_city.setdefault(text_key(row["cidade"]), set()).add(uf)
        inserted_or_seen = 0
        corrected = 0
        for row in rows:
            original_uf = None
            fixed_uf = normalize_uf(row["uf"])
            city_key_value = text_key(row["cidade"])
            if is_unknown_uf(fixed_uf):
                known_candidates = known_by_city.get(city_key_value, set())
                source_candidates = source_by_city.get(city_key_value, set())
                if len(known_candidates) == 1:
                    fixed_uf = next(iter(known_candidates))
                    original_uf = "NI"
                elif len(source_candidates) == 1:
                    fixed_uf = next(iter(source_candidates))
                    original_uf = "NI"
            if original_uf:
                corrected += 1
            upsert_geocity(conn, row["cidade"], fixed_uf, fonte="base", uf_corrigida_de=original_uf)
            inserted_or_seen += 1
        GEO_SYNC_LAST_RUN = now_ts
        return {"totalSeen": inserted_or_seen, "ufNiCorrected": corrected, "skipped": False}


def load_geocity_coordinates(conn):
    sync_status = {"ok": True, "error": None}
    try:
        sync_geocidades(conn)
    except sqlite3.OperationalError as exc:
        sync_status = {"ok": False, "error": str(exc)}
    rows = conn.execute(
        """
        SELECT cidade, uf, lat, lon, fonte, precisa_coordenada, uf_corrigida_de
        FROM geocidades
        """
    ).fetchall()
    coords = {}
    by_city = {}
    for row in rows:
        item = dict(row)
        coords[city_key(item["cidade"], item["uf"])] = item
        if item["lat"] is not None and item["lon"] is not None and not is_unknown_uf(item["uf"]):
            by_city.setdefault(text_key(item["cidade"]), []).append(item)
    return coords, by_city, sync_status


def geocity_lookup_rows(conn):
    rows = conn.execute(
        """
        SELECT cidade, uf, lat, lon, fonte, precisa_coordenada, uf_corrigida_de
        FROM geocidades
        """
    ).fetchall()
    coords = {}
    by_city = {}
    for row in rows:
        item = dict(row)
        coords[city_key(item["cidade"], item["uf"])] = item
        if item["lat"] is not None and item["lon"] is not None and not is_unknown_uf(item["uf"]):
            by_city.setdefault(text_key(item["cidade"]), []).append(item)
    return coords, by_city


def aggregate_geo_sources(conn):
    complaint_rows = conn.execute(
        """
        SELECT cidade, uf, COUNT(*) AS reclamacoes,
               SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
               0 AS assistencias,
               'complaints' AS origem
        FROM complaints
        WHERE cidade IS NOT NULL AND cidade != 'Sem cidade'
        GROUP BY cidade, uf
        """
    ).fetchall()
    operation_rows = conn.execute(
        """
        SELECT cidade, estado AS uf, 0 AS reclamacoes, 0 AS procedentes,
               COUNT(*) AS assistencias, 'operations' AS origem
        FROM operations
        WHERE cidade IS NOT NULL AND cidade != 'Sem cidade'
        GROUP BY cidade, estado
        """
    ).fetchall()
    rows = [dict(row) for row in complaint_rows] + [dict(row) for row in operation_rows]
    nps_exists = conn.execute(
        "SELECT COUNT(*) AS total FROM sqlite_master WHERE type='table' AND name='nps'"
    ).fetchone()["total"]
    if nps_exists:
        try:
            nps_rows = conn.execute(
                """
                SELECT cidade, uf, 0 AS reclamacoes, 0 AS procedentes,
                       COUNT(*) AS assistencias, 'nps' AS origem
                FROM nps
                WHERE cidade IS NOT NULL
                GROUP BY cidade, uf
                """
            ).fetchall()
            rows.extend(dict(row) for row in nps_rows)
        except sqlite3.Error:
            pass
    return rows


def atualizar_geocidades(force=True):
    init_db()
    with connect() as conn:
        sync_result = sync_geocidades(conn, force=force)
        conn.execute("DELETE FROM geocidades_pendentes WHERE ignorada = 0")
        coords, coords_by_city = geocity_lookup_rows(conn)
        processed = 0
        mapped_keys = set()
        pending_keys = set()
        for row in aggregate_geo_sources(conn):
            fixed_uf, original_uf, coord = resolve_geocity(coords, coords_by_city, row["cidade"], row["uf"])
            key = city_key(row["cidade"], fixed_uf)
            processed += 1
            if coord and coord.get("lat") is not None and coord.get("lon") is not None:
                mapped_keys.add(key)
                remove_pending_geocity(conn, row["cidade"], fixed_uf)
                continue
            pending_keys.add(key)
            upsert_pending_geocity(
                conn,
                row["cidade"],
                fixed_uf,
                original_uf or row["uf"],
                row["origem"],
                row.get("reclamacoes", 0),
                row.get("procedentes", 0),
                row.get("assistencias", 0),
            )
        total = len(mapped_keys | pending_keys)
        mapped = len(mapped_keys)
        pending = len(pending_keys)
        coverage = round(mapped / total, 4) if total else 1.0
        return {
            "processadas": processed,
            "mapeadas": mapped,
            "pendentes": pending,
            "cobertura": coverage,
            "coberturaPercentual": round(coverage * 100, 1),
            "sync": sync_result,
        }


def api_geo_coverage_snapshot():
    init_db()
    with connect() as conn:
        pending = conn.execute(
            """
            SELECT COUNT(*) AS value FROM (
                SELECT cidade_normalizada, uf
                FROM geocidades_pendentes
                WHERE ignorada = 0
                GROUP BY cidade_normalizada, uf
            )
            """
        ).fetchone()["value"] or 0
        mapped = conn.execute(
            """
            SELECT COUNT(*) AS value FROM (
                SELECT cidade_key, uf
                FROM geocidades
                WHERE lat IS NOT NULL AND lon IS NOT NULL
                GROUP BY cidade_key, uf
            )
            """
        ).fetchone()["value"] or 0
        total = mapped + pending
    missing = max(total - mapped, 0)
    coverage_pct = round(mapped * 100.0 / total, 1) if total else 100.0
    return {
        "mappedCities": mapped,
        "missingCoordinates": missing,
        "totalCities": total,
        "registeredCities": total,
        "registryCoveragePct": 100.0 if total else 100.0,
        "coveragePct": coverage_pct,
        "warning": coverage_pct < 95,
        "targetPct": 95,
    }


def api_geocidades_pendentes():
    init_db()
    with connect() as conn:
        rows = conn.execute(
            """
            SELECT cidade, uf, uf_original, origem, ocorrencias, reclamacoes,
                   procedentes, assistencias, created_at, updated_at
            FROM geocidades_pendentes
            WHERE ignorada = 0
            ORDER BY ocorrencias DESC, cidade ASC, uf ASC
            LIMIT 500
            """
        ).fetchall()
    return {"pendentes": [dict(row) for row in rows]}


def api_save_geocidade(payload):
    cidade = payload.get("cidade")
    uf = payload.get("uf")
    lat = payload.get("lat")
    lon = payload.get("lon")
    fonte = payload.get("fonte") or "manual"
    if not cidade or not uf:
        raise ValueError("Cidade e UF sao obrigatorias.")
    try:
        lat = float(lat)
        lon = float(lon)
    except (TypeError, ValueError) as exc:
        raise ValueError("Latitude e longitude devem ser numericas.") from exc
    if not (-34.5 <= lat <= 6.0 and -74.5 <= lon <= -32.0):
        raise ValueError("Coordenadas fora do intervalo esperado para o Brasil.")
    init_db()
    with connect() as conn:
        upsert_geocity(conn, cidade, uf, lat=lat, lon=lon, fonte=fonte)
        remove_pending_geocity(conn, cidade, uf)
    from cache_manager import clear_api_cache
    clear_api_cache()
    return {"ok": True, "geocidade": {"cidade": normalizar_cidade(cidade), "uf": normalize_uf(uf), "lat": lat, "lon": lon, "fonte": fonte}}


def api_ignore_geocidade(payload):
    cidade = payload.get("cidade")
    uf = normalize_uf(payload.get("uf"))
    if not cidade:
        raise ValueError("Cidade obrigatoria.")
    init_db()
    with connect() as conn:
        conn.execute(
            "UPDATE geocidades_pendentes SET ignorada = 1, updated_at = ? WHERE cidade_normalizada = ? AND uf = ?",
            (datetime.now().strftime("%Y-%m-%d %H:%M:%S"), normalizar_cidade(cidade), uf),
        )
    from cache_manager import clear_api_cache
    clear_api_cache()
    return {"ok": True}


def export_geocidades_pendentes_csv():
    payload = api_geocidades_pendentes()
    output = io.StringIO()
    fieldnames = ["cidade", "uf", "uf_original", "origem", "ocorrencias", "reclamacoes", "procedentes", "assistencias"]
    writer = csv.DictWriter(output, fieldnames=fieldnames, extrasaction="ignore")
    writer.writeheader()
    writer.writerows(payload["pendentes"])
    return output.getvalue().encode("utf-8-sig")


def resolve_geocity(coords, by_city, cidade, uf):
    fixed_uf = normalize_uf(uf)
    original_uf = None
    key = city_key(cidade, fixed_uf)
    coord = coords.get(key)
    if coord and coord.get("lat") is not None and coord.get("lon") is not None:
        return fixed_uf, original_uf, coord
    candidates = by_city.get(text_key(cidade), [])
    if is_unknown_uf(fixed_uf) and len(candidates) == 1:
        coord = candidates[0]
        return normalize_uf(coord["uf"]), original_uf or normalize_uf(uf), coord
    return fixed_uf, original_uf, coord


def map_status(row):
    procedencia = row["taxa_procedencia"] or 0
    perda = row["taxa_perda"] or 0
    sucesso = row["taxa_sucesso"] if row["taxa_sucesso"] is not None else 100
    if procedencia >= 65 or perda >= 35 or sucesso < 60:
        return "critico"
    if procedencia >= 45 or perda >= 20 or sucesso < 75:
        return "atencao"
    return "saudavel"


def api_map(query):
    init_db()
    params = {key: values[0] for key, values in query.items() if values and values[0]}
    with connect() as conn:
        coords, coords_by_city, sync_status = load_geocity_coordinates(conn)
        max_date = conn.execute("SELECT MAX(data_reclamacao) AS value FROM complaints").fetchone()["value"]
        max_operation_date = conn.execute("SELECT MAX(data) AS value FROM operations").fetchone()["value"]
        params["_max_date"] = max_date
        params["_max_operation_date"] = max_operation_date
        complaint_where, complaint_args = build_where(params, procedentes_only=False)
        operation_where, operation_args = build_operation_where(params)
        complaint_rows = conn.execute(
            f"""
            SELECT cidade, uf,
                   COUNT(*) AS reclamacoes,
                   SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedencias
            FROM complaints
            {complaint_where}
            GROUP BY cidade, uf
            """,
            complaint_args,
        ).fetchall()
        top_hub_rows = conn.execute(
            f"""
            SELECT cidade, uf, hub AS value, COUNT(*) AS total
            FROM complaints
            {complaint_where}
            GROUP BY cidade, uf, hub
            ORDER BY cidade, uf, total DESC, hub ASC
            """,
            complaint_args,
        ).fetchall()
        top_client_rows = conn.execute(
            f"""
            SELECT cidade, uf, cliente AS value, COUNT(*) AS total
            FROM complaints
            {complaint_where}
            GROUP BY cidade, uf, cliente
            ORDER BY cidade, uf, total DESC, cliente ASC
            """,
            complaint_args,
        ).fetchall()
        top_reason_rows = conn.execute(
            f"""
            SELECT cidade, uf, motivo AS value, COUNT(*) AS total
            FROM complaints
            {complaint_where + (" AND procedente = 1" if complaint_where else " WHERE procedente = 1")}
            GROUP BY cidade, uf, motivo
            ORDER BY cidade, uf, total DESC, motivo ASC
            """,
            complaint_args,
        ).fetchall()
        operation_rows = conn.execute(
            f"""
            SELECT cidade, estado AS uf,
                   COUNT(*) AS assistencias,
                   SUM(frustrada) AS frustracoes,
                   SUM(perda) AS perdas,
                   SUM(concluida) AS concluidas
            FROM operations
            {operation_where}
            GROUP BY cidade, estado
            """,
            operation_args,
        ).fetchall()
        top_operation_reason_rows = conn.execute(
            f"""
            SELECT cidade, estado AS uf, motivo_nao_realizada AS value, COUNT(*) AS total
            FROM operations
            {operation_where}
            GROUP BY cidade, estado, motivo_nao_realizada
            ORDER BY cidade, estado, total DESC, motivo_nao_realizada ASC
            """,
            operation_args,
        ).fetchall()

        def top_map(rows):
            output = {}
            for item in rows:
                key = city_key(item["cidade"], item["uf"])
                if key not in output and item["value"]:
                    output[key] = item["value"]
            return output

        top_hub = top_map(top_hub_rows)
        top_client = top_map(top_client_rows)
        top_reason = top_map(top_reason_rows)
        top_operation_reason = top_map(top_operation_reason_rows)

        merged = {}

        def ensure_city(row):
            fixed_uf, original_uf, coord = resolve_geocity(coords, coords_by_city, row["cidade"], row["uf"])
            key = city_key(row["cidade"], fixed_uf)
            merged.setdefault(
                key,
                {
                    "cidade": row["cidade"],
                    "uf": fixed_uf,
                    "uf_original": original_uf,
                    "reclamacoes": 0,
                    "procedencias": 0,
                    "assistencias": 0,
                    "frustracoes": 0,
                    "perdas": 0,
                    "concluidas": 0,
                    "principal_hub": None,
                    "principal_cliente": None,
                    "principal_motivo": None,
                    "principal_motivo_operacional": None,
                    "_coord": coord,
                },
            )
            if not merged[key].get("_coord") and coord:
                merged[key]["_coord"] = coord
            return merged[key]

        for row in complaint_rows:
            item = ensure_city(row)
            item["reclamacoes"] += row["reclamacoes"] or 0
            item["procedencias"] += row["procedencias"] or 0
            source_key = city_key(row["cidade"], row["uf"])
            item["principal_hub"] = top_hub.get(source_key) or item["principal_hub"]
            item["principal_motivo"] = top_reason.get(source_key) or item["principal_motivo"]
            item["principal_cliente"] = top_client.get(source_key) or item["principal_cliente"]
        for row in operation_rows:
            item = ensure_city(row)
            item["assistencias"] += row["assistencias"] or 0
            item["frustracoes"] += row["frustracoes"] or 0
            item["perdas"] += row["perdas"] or 0
            item["concluidas"] += row["concluidas"] or 0
            item["principal_motivo_operacional"] = top_operation_reason.get(city_key(row["cidade"], row["uf"])) or item["principal_motivo_operacional"]

    points = []
    missing_cities = []
    for key, item in merged.items():
        coord = item.pop("_coord", None)
        if not coord or coord.get("lat") is None or coord.get("lon") is None:
            missing_cities.append(
                {
                    "cidade": item["cidade"],
                    "uf": item["uf"],
                    "uf_original": item.get("uf_original"),
                    "reclamacoes": item["reclamacoes"],
                    "procedencias": item["procedencias"],
                    "assistencias": item["assistencias"],
                }
            )
            continue
        item["lat"] = coord["lat"]
        item["lon"] = coord["lon"]
        item["taxa_procedencia"] = round(item["procedencias"] * 100.0 / item["reclamacoes"], 1) if item["reclamacoes"] else 0
        item["taxa_perda"] = round(item["perdas"] * 100.0 / item["assistencias"], 1) if item["assistencias"] else 0
        item["taxa_sucesso"] = round(item["concluidas"] * 100.0 / item["assistencias"], 1) if item["assistencias"] else None
        item["status"] = map_status(item)
        item["volume_total"] = item["reclamacoes"] + item["assistencias"]
        points.append(item)
    points.sort(key=lambda row: row["volume_total"], reverse=True)
    missing_cities.sort(key=lambda row: row["reclamacoes"] + row["assistencias"], reverse=True)
    total_cities = len(points) + len(missing_cities)
    coverage_pct = round(len(points) * 100.0 / total_cities, 1) if total_cities else 100.0
    return {
        "points": points,
        "mappedCities": len(points),
        "missingCoordinates": len(missing_cities),
        "totalCities": total_cities,
        "registeredCities": total_cities,
        "registryCoveragePct": 100.0 if total_cities else 100.0,
        "coveragePct": coverage_pct,
        "coverageWarning": coverage_pct < 95,
        "geoCoverage": {
            "mappedCities": len(points),
            "missingCoordinates": len(missing_cities),
            "totalCities": total_cities,
            "registeredCities": total_cities,
            "registryCoveragePct": 100.0 if total_cities else 100.0,
            "coveragePct": coverage_pct,
            "warning": coverage_pct < 95,
            "targetPct": 95,
        },
        "missingCities": missing_cities[:100],
        "citiesFile": str(CITIES_PATH),
        "source": "geocidades",
        "geoSync": sync_status,
    }


def api_providers(query):
    init_db()
    params = {key: values[0] for key, values in query.items() if values and values[0]}
    provider_filter = params.pop("prestador", None)
    with connect() as conn:
        max_date = conn.execute("SELECT MAX(data_reclamacao) AS value FROM complaints").fetchone()["value"]
        params["_max_date"] = max_date
        where, args = build_where(params, procedentes_only=False)
        if provider_filter:
            where = where + (" AND prestador = ?" if where else " WHERE prestador = ?")
            args.append(provider_filter)
        rows = conn.execute(
            f"""
            WITH base AS (
              SELECT
                prestador,
                hub,
                regiao,
                cidade,
                cliente,
                mes,
                procedente,
                validacao,
                sla_ok,
                dias_uteis_fechamento,
                status_reclamacao
              FROM complaints
              {where}
            ),
            agg AS (
              SELECT
                prestador,
                hub,
                regiao,
                COUNT(*) AS reclamacoes,
                SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes,
                SUM(CASE WHEN UPPER(validacao) = 'IMPROCEDENTE' THEN 1 ELSE 0 END) AS improcedentes,
                SUM(CASE WHEN sla_ok = 0 AND dias_uteis_fechamento IS NOT NULL THEN 1 ELSE 0 END) AS sla_estourado,
                COUNT(DISTINCT cidade) AS cidades,
                ROUND(SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) * 100.0 / NULLIF(COUNT(*), 0), 1) AS taxa_procedencia,
                ROUND(SUM(CASE WHEN sla_ok = 1 THEN 1 ELSE 0 END) * 100.0 / NULLIF(SUM(CASE WHEN UPPER(status_reclamacao) = 'FECHADA' THEN 1 ELSE 0 END), 0), 1) AS taxa_sla,
                ROUND(AVG(CASE WHEN dias_uteis_fechamento IS NOT NULL THEN dias_uteis_fechamento END), 2) AS media_dias_uteis
              FROM base
              GROUP BY prestador, hub, regiao
            ),
            maxes AS (
              SELECT MAX(reclamacoes) AS max_reclamacoes FROM agg
            )
            SELECT
              agg.*,
              ROUND(
                40.0 * reclamacoes / NULLIF(max_reclamacoes, 0) +
                30.0 * COALESCE(taxa_procedencia, 0) / 100.0 +
                20.0 * (1.0 - COALESCE(taxa_sla, 0) / 100.0) +
                10.0 * CASE WHEN cidades > 1 THEN 1 ELSE 0 END
              , 1) AS score_risco
            FROM agg, maxes
            ORDER BY score_risco DESC, reclamacoes DESC, prestador ASC
            LIMIT 120
            """,
            args,
        ).fetchall()
        providers = []
        for row in rows:
            item = dict(row)
            item["status_risco"] = risk_status(item["score_risco"] or 0)
            providers.append(item)
        detail = None
        if provider_filter:
            monthly = conn.execute(
                f"""
                SELECT mes AS label, COUNT(*) AS reclamacoes,
                       SUM(CASE WHEN procedente = 1 THEN 1 ELSE 0 END) AS procedentes
                FROM complaints
                {where + (" AND mes != 'Sem data'" if where else " WHERE mes != 'Sem data'")}
                GROUP BY mes
                ORDER BY mes ASC
                """,
                args,
            ).fetchall()
            cities = conn.execute(
                f"""
                SELECT cidade AS label, COUNT(*) AS value
                FROM complaints
                {where}
                GROUP BY cidade
                ORDER BY value DESC, cidade ASC
                LIMIT 20
                """,
                args,
            ).fetchall()
            records = conn.execute(
                f"""
                SELECT data_reclamacao, dt_fechamento, cliente, cidade, regiao, hub, status_reclamacao,
                       validacao, motivo, analista_sac, dias_uteis_fechamento, sla_ok
                FROM complaints
                {where}
                ORDER BY COALESCE(dt_fechamento, data_reclamacao) DESC
                LIMIT 100
                """,
                args,
            ).fetchall()
            detail = {
                "provider": provider_filter,
                "monthly": [dict(row) for row in monthly],
                "cities": [dict(row) for row in cities],
                "records": [dict(row) for row in records],
            }
        return {"providers": providers, "detail": detail}


def api_data_freshness():
    init_db()
    with connect() as conn:
        row = conn.execute(
            """
            SELECT
              (SELECT MAX(data_reclamacao) FROM complaints) AS max_complaint_date,
              (SELECT MAX(data) FROM operations) AS max_operation_date,
              (SELECT COUNT(*) FROM complaints) AS complaints_total,
              (SELECT COUNT(*) FROM operations) AS operations_total
            """
        ).fetchone()
    complaint_date = row["max_complaint_date"]
    operation_date = row["max_operation_date"]
    lag_days = None
    if complaint_date and operation_date:
        complaint_ts = pd.to_datetime(complaint_date, errors="coerce")
        operation_ts = pd.to_datetime(operation_date, errors="coerce")
        if not pd.isna(complaint_ts) and not pd.isna(operation_ts):
            lag_days = abs((complaint_ts.date() - operation_ts.date()).days)
    return {
        "maxComplaintDate": complaint_date,
        "maxOperationDate": operation_date,
        "lagDays": lag_days,
        "warningDays": DATE_LAG_WARNING_DAYS,
        "complaintsTotal": row["complaints_total"] or 0,
        "operationsTotal": row["operations_total"] or 0,
    }


def _pct_value(value):
    return round((value or 0) * 100, 1)


def _alert_priority(severity, impact=3, urgency=3):
    severity_weight = {"Vermelho": 3, "Amarelo": 2, "Verde": 1}.get(severity, 1)
    return round((severity_weight * 30) + (impact * 5) + (urgency * 5), 1)


def _make_alert(
    severity,
    area,
    title,
    message,
    value=None,
    recommendation=None,
    action=None,
    driver=None,
    dimension=None,
    impact=3,
    urgency=3,
):
    if recommendation is None:
        if severity == "Vermelho":
            recommendation = "Priorizar plano de acao e acompanhar diariamente."
        elif severity == "Amarelo":
            recommendation = "Acompanhar tendencia e atuar antes de virar critico."
        else:
            recommendation = "Manter acompanhamento operacional."
    return {
        "severity": severity,
        "area": area,
        "message": message,
        "titulo": title,
        "descricao": message,
        "dimensao": dimension or area,
        "valor": value,
        "driver": driver,
        "impacto": impact,
        "urgencia": urgency,
        "prioridade": _alert_priority(severity, impact, urgency),
        "recomendacao": recommendation,
        "acao": action or recommendation,
    }


def _share_alert(rows, total, area, title_prefix, red_threshold, yellow_threshold, action):
    if not rows or not total:
        return None
    top = rows[0]
    label = top.get("label") or top.get("cidade") or top.get("cliente") or top.get("motivo") or "Sem identificacao"
    value = top.get("value", 0)
    share = value / total if total else 0
    if share >= red_threshold:
        severity = "Vermelho"
    elif share >= yellow_threshold:
        severity = "Amarelo"
    else:
        return None
    pct_share = _pct_value(share)
    return _make_alert(
        severity,
        area,
        f"{title_prefix} concentrado",
        f"{label} concentra {pct_share}% das procedencias do recorte.",
        value=pct_share,
        driver=label,
        action=action,
        recommendation="Quebrar o indicador por HUB, prestador e motivo para definir responsavel e plano de reducao.",
        impact=4 if severity == "Vermelho" else 3,
        urgency=4 if severity == "Vermelho" else 3,
    )


def api_alerts(query):
    summary = api_summary(query)
    productivity = api_productivity(query)
    operation = api_operation(query)
    providers = api_providers(query)["providers"]
    geo_coverage = api_geo_coverage_snapshot()
    freshness = api_data_freshness()
    alerts = []
    if (
        freshness["lagDays"] is not None
        and freshness["lagDays"] > freshness["warningDays"]
        and freshness["complaintsTotal"]
        and freshness["operationsTotal"]
    ):
        alerts.append(
            _make_alert(
                "Amarelo",
                "Bases",
                "Bases fora de sincronia",
                (
                    f"Reclamacoes ate {freshness['maxComplaintDate']} e operacao ate "
                    f"{freshness['maxOperationDate']} ({freshness['lagDays']} dias de diferenca)."
                ),
                value=freshness["lagDays"],
                action="Atualizar as duas bases antes de comparar reclamacoes com operacao.",
                driver="Diferenca entre datas maximas das bases",
                impact=3,
                urgency=2,
            )
        )
    comp = summary["kpis"]["comparison"]
    if comp["previousMonth"] and comp["deltaPct"] > 0.20:
        alerts.append(
            _make_alert(
                "Vermelho",
                "Reclamacoes",
                "Alta mensal de procedentes",
                f"Procedentes cresceram {round(comp['deltaPct'] * 100, 1)}% contra o mes anterior.",
                value=round(comp["deltaPct"] * 100, 1),
                action="Comparar os motivos e HUBs do mes atual contra o mes anterior.",
                driver="Mes atual x mes anterior",
                impact=5,
                urgency=4,
            )
        )
    if summary["kpis"]["procedenciaPct"] > 0.65:
        alerts.append(
            _make_alert(
                "Vermelho",
                "Qualidade",
                "Procedencia acima do limite critico",
                f"Taxa de procedencia em {_pct_value(summary['kpis']['procedenciaPct'])}%, acima do limite critico.",
                value=_pct_value(summary["kpis"]["procedenciaPct"]),
                action="Atuar nos tres maiores motivos procedentes e revisar causa raiz por HUB.",
                driver="Taxa de procedencia",
                impact=5,
                urgency=5,
            )
        )
    elif summary["kpis"]["procedenciaPct"] > 0.45:
        alerts.append(
            _make_alert(
                "Amarelo",
                "Qualidade",
                "Procedencia em atencao",
                f"Taxa de procedencia em {_pct_value(summary['kpis']['procedenciaPct'])}%, exige acompanhamento.",
                value=_pct_value(summary["kpis"]["procedenciaPct"]),
                action="Monitorar a evolucao semanal e atacar motivos com crescimento.",
                driver="Taxa de procedencia",
            )
        )
    if productivity["kpis"]["slaPct"] < 0.50:
        alerts.append(
            _make_alert(
                "Vermelho",
                "SLA",
                "SLA SAC critico",
                f"SLA SAC em {_pct_value(productivity['kpis']['slaPct'])}%, abaixo da meta de fechamento em ate {SLA_BUSINESS_DAYS_TARGET} dias uteis.",
                value=_pct_value(productivity["kpis"]["slaPct"]),
                action="Priorizar analistas com menor taxa de SLA e backlog vencido.",
                driver="Tempo entre abertura e fechamento",
                impact=5,
                urgency=5,
            )
        )
    elif productivity["kpis"]["slaPct"] < 0.75:
        alerts.append(
            _make_alert(
                "Amarelo",
                "SLA",
                "SLA SAC em atencao",
                f"SLA SAC em {_pct_value(productivity['kpis']['slaPct'])}%, em atencao.",
                value=_pct_value(productivity["kpis"]["slaPct"]),
                action="Acompanhar fechamentos diarios e redistribuir fila quando necessario.",
                driver="Tempo entre abertura e fechamento",
                impact=4,
            )
        )
    backlog = productivity["kpis"].get("backlog", 0) or 0
    total_productivity = productivity["kpis"].get("total", 0) or 0
    backlog_share = backlog / total_productivity if total_productivity else 0
    if backlog and backlog_share >= 0.30:
        alerts.append(
            _make_alert(
                "Vermelho",
                "Backlog",
                "Backlog elevado",
                f"Backlog representa {_pct_value(backlog_share)}% das reclamacoes do recorte ({backlog} registros).",
                value=_pct_value(backlog_share),
                action="Criar fila priorizada por idade e analista responsavel.",
                driver="Reclamacoes ainda nao fechadas",
                impact=5,
                urgency=5,
            )
        )
    elif backlog and backlog_share >= 0.15:
        alerts.append(
            _make_alert(
                "Amarelo",
                "Backlog",
                "Backlog em observacao",
                f"Backlog representa {_pct_value(backlog_share)}% das reclamacoes do recorte ({backlog} registros).",
                value=_pct_value(backlog_share),
                action="Verificar analistas com maior fila aberta e registros acima de 5 dias uteis.",
                driver="Reclamacoes ainda nao fechadas",
                impact=3,
            )
        )
    share_alerts = [
        _share_alert(
            summary["charts"].get("cidade", []),
            summary["kpis"].get("procedentes", 0),
            "Cidade",
            "Cidade",
            0.30,
            0.20,
            "Abrir ranking por cidade e cruzar com HUB, prestador e motivo.",
        ),
        _share_alert(
            summary["charts"].get("cliente", []),
            summary["kpis"].get("procedentes", 0),
            "Cliente",
            "Cliente corporativo",
            0.30,
            0.20,
            "Validar se o cliente possui padrao especifico de atendimento ou recorrencia operacional.",
        ),
        _share_alert(
            summary["charts"].get("motivo", []),
            summary["kpis"].get("procedentes", 0),
            "Motivo",
            "Motivo",
            0.35,
            0.25,
            "Transformar o motivo em plano de causa raiz com dono e prazo.",
        ),
    ]
    alerts.extend([item for item in share_alerts if item])
    if operation["kpis"].get("successPct", 0) < 0.60 and operation["kpis"].get("total", 0):
        alerts.append(
            _make_alert(
                "Vermelho",
                "Operacao",
                "Sucesso operacional critico",
                f"Taxa de sucesso operacional em {_pct_value(operation['kpis'].get('successPct', 0))}%.",
                value=_pct_value(operation["kpis"].get("successPct", 0)),
                action="Priorizar HUBs com maior perda e motivos de frustracao/cancelamento.",
                driver="Assistencias concluidas sobre total",
                impact=5,
                urgency=4,
            )
        )
    elif operation["kpis"].get("successPct", 0) < 0.75 and operation["kpis"].get("total", 0):
        alerts.append(
            _make_alert(
                "Amarelo",
                "Operacao",
                "Sucesso operacional em atencao",
                f"Taxa de sucesso operacional em {_pct_value(operation['kpis'].get('successPct', 0))}%.",
                value=_pct_value(operation["kpis"].get("successPct", 0)),
                action="Monitorar perdas por HUB e motivo para reduzir reentregas e improdutividade.",
                driver="Assistencias concluidas sobre total",
                impact=4,
                urgency=3,
            )
        )
    if providers:
        top = providers[0]
        if top["score_risco"] >= 70:
            alerts.append(
                _make_alert(
                    "Vermelho",
                    "Prestador",
                    "Prestador em risco critico",
                    f"Prestador {top['prestador']} com score de risco {top['score_risco']}.",
                    value=top["score_risco"],
                    action="Abrir detalhe do prestador e revisar ocorrencias por cidade, SLA e motivo.",
                    driver=top["prestador"],
                    impact=5,
                    urgency=4,
                )
            )
        elif top["score_risco"] >= 40:
            alerts.append(
                _make_alert(
                    "Amarelo",
                    "Prestador",
                    "Prestador em atencao",
                    f"Prestador {top['prestador']} em atencao, score {top['score_risco']}.",
                    value=top["score_risco"],
                    action="Acompanhar tendencia do prestador antes de evoluir para critico.",
                    driver=top["prestador"],
                    impact=3,
                )
            )
    worst_hub = operation["kpis"].get("worstHub")
    if worst_hub and (worst_hub.get("taxa_perda") or 0) >= 30:
        alerts.append(
            _make_alert(
                "Vermelho",
                "HUB",
                "HUB com perda operacional critica",
                f"HUB {worst_hub['label']} com {worst_hub['taxa_perda']}% de perda operacional.",
                value=worst_hub["taxa_perda"],
                action="Quebrar o HUB por montador/prestador e motivo de nao realizacao.",
                driver=worst_hub["label"],
                impact=5,
                urgency=4,
            )
        )
    missing_coordinates = geo_coverage.get("missingCoordinates", 0) or 0
    if missing_coordinates:
        alerts.append(
            _make_alert(
                "Amarelo",
                "Mapa",
                "Cobertura geografica abaixo da meta" if geo_coverage.get("warning") else "Cidades sem coordenadas",
                (
                    f"{missing_coordinates} cidades estao sem latitude/longitude. "
                    f"Cobertura geografica atual: {geo_coverage.get('coveragePct', 0)}%."
                ),
                value=missing_coordinates,
                action="Completar latitude e longitude na tabela geocidades para as cidades listadas no painel.",
                driver="Cobertura geografica",
                impact=3 if geo_coverage.get("warning") else 2,
                urgency=3 if geo_coverage.get("warning") else 2,
            )
        )
    if not alerts:
        alerts.append(
            _make_alert(
                "Verde",
                "Geral",
                "Operacao sem alertas criticos",
                "Nenhum alerta critico nos filtros atuais.",
                action="Manter acompanhamento e comparar com periodos anteriores.",
                impact=1,
                urgency=1,
            )
        )
    alerts = sorted(alerts, key=lambda item: item.get("prioridade", 0), reverse=True)
    return {"alerts": alerts}


def api_insights(query):
    summary = api_summary(query)
    productivity = api_productivity(query)
    operation = api_operation(query)
    providers = api_providers(query)["providers"]
    geo_coverage = api_geo_coverage_snapshot()
    insights = []
    cards = []
    top_hub = summary["charts"]["hub"][0] if summary["charts"]["hub"] else None
    top_reason = summary["charts"]["motivo"][0] if summary["charts"]["motivo"] else None
    top_region = summary["charts"]["regiao"][0] if summary["charts"]["regiao"] else None
    top_city = summary["charts"]["cidade"][0] if summary["charts"]["cidade"] else None
    top_client = summary["charts"]["cliente"][0] if summary["charts"]["cliente"] else None
    top_provider = providers[0] if providers else None
    procedentes = summary["kpis"].get("procedentes", 0) or 0

    def add_card(kind, title, description, action, impact="Medio", priority=50):
        cards.append(
            {
                "tipo": kind,
                "titulo": title,
                "descricao": description,
                "acao": action,
                "impacto": impact,
                "prioridade": priority,
            }
        )

    if top_hub and summary["kpis"]["procedentes"]:
        share = top_hub["value"] / summary["kpis"]["procedentes"]
        insights.append(f"O HUB {top_hub['label']} concentra {round(share * 100, 1)}% das procedencias do recorte.")
        add_card(
            "Concentracao",
            f"HUB {top_hub['label']} puxa procedencias",
            f"Ele concentra {round(share * 100, 1)}% das reclamacoes procedentes do recorte.",
            "Abrir o detalhamento por motivo, cidade e prestador para isolar causa raiz.",
            "Alto" if share >= 0.25 else "Medio",
            80 if share >= 0.25 else 60,
        )
    if top_reason and summary["kpis"]["procedentes"]:
        share = top_reason["value"] / summary["kpis"]["procedentes"]
        insights.append(f"O motivo {top_reason['label']} representa {round(share * 100, 1)}% das procedencias.")
        add_card(
            "Causa raiz",
            f"Motivo dominante: {top_reason['label']}",
            f"O motivo representa {round(share * 100, 1)}% das procedencias.",
            "Converter esse motivo em plano operacional com responsavel, prazo e acompanhamento semanal.",
            "Alto" if share >= 0.30 else "Medio",
            85 if share >= 0.30 else 65,
        )
    if top_region and procedentes:
        share = top_region["value"] / procedentes
        add_card(
            "Regional",
            f"Regional {top_region['label']} lidera volume",
            f"A regional concentra {round(share * 100, 1)}% das procedencias.",
            "Comparar HUBs da regional e priorizar o que combina alto volume com baixa taxa de SLA.",
            "Medio",
            62,
        )
    if top_city and procedentes:
        share = top_city["value"] / procedentes
        add_card(
            "Geografia",
            f"Cidade mais sensivel: {top_city['label']}",
            f"A cidade representa {round(share * 100, 1)}% das procedencias.",
            "Cruzar a cidade com prestador e HUB para identificar concentracao operacional.",
            "Alto" if share >= 0.20 else "Medio",
            78 if share >= 0.20 else 55,
        )
    if top_client and procedentes:
        share = top_client["value"] / procedentes
        add_card(
            "Cliente",
            f"Cliente corporativo em foco: {top_client['label']}",
            f"O cliente responde por {round(share * 100, 1)}% das procedencias.",
            "Verificar se ha padrao de produto, rota, agendamento ou expectativa de atendimento.",
            "Alto" if share >= 0.20 else "Medio",
            76 if share >= 0.20 else 54,
        )
    if operation["kpis"].get("worstHub"):
        hub = operation["kpis"]["worstHub"]
        insights.append(f"O HUB {hub['label']} tem a maior taxa de perda operacional: {hub['taxa_perda']}%.")
        add_card(
            "Operacao",
            f"Maior perda operacional: {hub['label']}",
            f"O HUB registra {hub['taxa_perda']}% de perda operacional.",
            "Abrir perdas por motivo e prestador para reduzir cancelamentos, frustracoes e improdutividade.",
            "Alto" if (hub.get("taxa_perda") or 0) >= 30 else "Medio",
            82 if (hub.get("taxa_perda") or 0) >= 30 else 58,
        )
    if productivity["kpis"].get("leadAnalyst"):
        lead = productivity["kpis"]["leadAnalyst"]
        insights.append(f"Maior volume SAC no recorte: {lead['label']} com {lead['value']} reclamacoes tratadas.")
        add_card(
            "Produtividade",
            f"Referencial de volume: {lead['label']}",
            f"O analista tratou {lead['value']} reclamacoes no recorte.",
            "Comparar volume com SLA e backlog para diferenciar alta produtividade de sobrecarga.",
            "Medio",
            52,
        )
    sla_pct = productivity["kpis"].get("slaPct", 0) or 0
    if productivity["kpis"].get("slaTotal", 0):
        add_card(
            "SLA",
            f"SLA de fechamento em {_pct_value(sla_pct)}%",
            f"A meta considera fechamento em ate {SLA_BUSINESS_DAYS_TARGET} dias uteis, sem finais de semana.",
            "Priorizar analistas com maior media de dias uteis e maior backlog vencido.",
            "Alto" if sla_pct < 0.75 else "Baixo",
            88 if sla_pct < 0.50 else 70 if sla_pct < 0.75 else 35,
        )
    backlog = productivity["kpis"].get("backlog", 0) or 0
    if backlog:
        add_card(
            "Backlog",
            f"{backlog} reclamacoes em aberto",
            "O backlog e o principal indicador antecipado de risco para o SLA dos proximos dias.",
            "Criar rotina diaria de fila por idade, analista e cliente corporativo.",
            "Alto" if backlog >= 100 else "Medio",
            84 if backlog >= 100 else 57,
        )
    if top_provider:
        insights.append(f"Prestador com maior risco: {top_provider['prestador']} com score {top_provider['score_risco']} ({top_provider['status_risco']}).")
        add_card(
            "Prestador",
            f"Prestador com maior score: {top_provider['prestador']}",
            f"Score de risco {top_provider['score_risco']} ({top_provider['status_risco']}).",
            "Abrir detalhe do prestador e revisar cidades, SLA e motivos associados.",
            "Alto" if top_provider["score_risco"] >= 70 else "Medio",
            86 if top_provider["score_risco"] >= 70 else 60,
        )
    missing_coordinates = geo_coverage.get("missingCoordinates", 0) or 0
    if missing_coordinates:
        add_card(
            "Dados",
            f"Cobertura geografica em {geo_coverage.get('coveragePct', 0)}%",
            f"{missing_coordinates} cidades entram nos rankings, mas ainda nao aparecem visualmente no mapa geografico.",
            "Completar coordenadas na tabela geocidades; novas cidades ja sao cadastradas automaticamente.",
            "Medio" if geo_coverage.get("warning") else "Baixo",
            68 if geo_coverage.get("warning") else 40,
        )
    cards = sorted(cards, key=lambda item: item.get("prioridade", 0), reverse=True)
    return {"insights": insights, "cards": cards[:12]}


def export_executive_html(query):
    summary = api_summary(query)
    productivity = api_productivity(query)
    operation = api_operation(query)
    alerts = api_alerts(query)["alerts"]
    html = f"""<!doctype html><html><head><meta charset='utf-8'><title>Resumo Executivo MMS</title>
    <style>body{{font-family:Calibri,Arial;margin:32px;color:#171a45}}h1{{color:#2B3594}}.card{{border:1px solid #d8dbee;padding:12px;margin:10px 0;border-left:6px solid #FF8C28}}</style></head><body>
    <h1>Resumo Executivo MMS</h1>
    <div class='card'>Reclamacoes: {summary['kpis']['total']} | Procedentes: {summary['kpis']['procedentes']} | Taxa: {round(summary['kpis']['procedenciaPct']*100,1)}%</div>
    <div class='card'>SLA SAC: {round(productivity['kpis']['slaPct']*100,1)}% | Fechadas: {productivity['kpis']['fechadas']} | Backlog: {productivity['kpis']['backlog']}</div>
    <div class='card'>Operacao: {operation['kpis']['total']} assistencias | Sucesso: {round(operation['kpis']['successPct']*100,1)}% | Perdas: {operation['kpis']['perdas']}</div>
    <h2>Alertas</h2>{''.join(f"<div class='card'><strong>{a['severity']} - {a['area']}</strong><br>{a['message']}</div>" for a in alerts)}
    <script>window.print()</script></body></html>"""
    return html.encode("utf-8-sig")


def export_excel_workbook(query):
    output = io.BytesIO()
    summary = api_summary(query)
    productivity = api_productivity(query)
    operation = api_operation(query)
    alerts = api_alerts(query)["alerts"]
    providers = api_providers(query)["providers"]
    scores = {
        "Prestadores": api_score_dimension(query, "prestadores")["rows"],
        "HUBs": api_score_dimension(query, "hubs")["rows"],
        "Cidades": api_score_dimension(query, "cidades")["rows"],
        "Clientes": api_score_dimension(query, "clientes")["rows"],
    }
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        pd.DataFrame([
            {"Indicador": "Reclamacoes", "Valor": summary["kpis"]["total"]},
            {"Indicador": "Procedentes", "Valor": summary["kpis"]["procedentes"]},
            {"Indicador": "Taxa procedencia", "Valor": round(summary["kpis"]["procedenciaPct"] * 100, 1)},
            {"Indicador": "SLA SAC", "Valor": round(productivity["kpis"]["slaPct"] * 100, 1)},
            {"Indicador": "Sucesso operacional", "Valor": round(operation["kpis"]["successPct"] * 100, 1)},
        ]).to_excel(writer, sheet_name="Resumo Executivo", index=False)
        pd.DataFrame(summary["tables"]["combo"]).to_excel(writer, sheet_name="Reclamacoes", index=False)
        pd.DataFrame(productivity["tables"]["analysts"]).to_excel(writer, sheet_name="Produtividade", index=False)
        pd.DataFrame(operation["tables"]["hubPerformance"]).to_excel(writer, sheet_name="Operacao", index=False)
        pd.DataFrame(alerts).to_excel(writer, sheet_name="Alertas", index=False)
        start_row = 0
        for name, rows in scores.items():
            pd.DataFrame(rows).to_excel(writer, sheet_name="Scores", startrow=start_row, index=False)
            start_row += max(2, len(rows) + 3)
        pd.DataFrame(providers).to_excel(writer, sheet_name="Prestadores", index=False)
    return output.getvalue()


def export_executive_pdf(query):
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    output = io.BytesIO()
    summary = api_summary(query)
    productivity = api_productivity(query)
    operation = api_operation(query)
    alerts = api_alerts(query)["alerts"]
    doc = SimpleDocTemplate(output, pagesize=A4, title="Resumo Executivo MMS")
    styles = getSampleStyleSheet()
    story = [
        Paragraph("Resumo Executivo MMS", styles["Title"]),
        Paragraph("Central de Performance Operacional", styles["Normal"]),
        Spacer(1, 12),
    ]
    kpis = [
        ["Indicador", "Valor"],
        ["Reclamacoes", summary["kpis"]["total"]],
        ["Procedentes", summary["kpis"]["procedentes"]],
        ["Taxa procedencia", f"{round(summary['kpis']['procedenciaPct'] * 100, 1)}%"],
        ["SLA SAC", f"{round(productivity['kpis']['slaPct'] * 100, 1)}%"],
        ["Sucesso operacional", f"{round(operation['kpis']['successPct'] * 100, 1)}%"],
        ["Perdas operacionais", operation["kpis"]["perdas"]],
    ]
    table = Table(kpis, hAlign="LEFT")
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2B3594")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#d8dbee")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("PADDING", (0, 0), (-1, -1), 6),
    ]))
    story += [table, Spacer(1, 16), Paragraph("Alertas", styles["Heading2"])]
    for alert in alerts[:10]:
        story.append(Paragraph(f"<b>{alert['severity']} - {alert['area']}</b>: {alert['message']}", styles["Normal"]))
        story.append(Spacer(1, 6))
    doc.build(story)
    return output.getvalue()


def export_executive_pptx(query):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    output = io.BytesIO()
    summary = api_summary(query)
    productivity = api_productivity(query)
    operation = api_operation(query)
    alerts = api_alerts(query)["alerts"]
    insights = api_insights(query)["insights"]
    map_payload = api_map(query)
    prs = Presentation()

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_bullets(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.font.size = Pt(18)

    add_title_slide("Dashboard MMS", "Resumo executivo operacional")
    add_bullets("Resumo Executivo", [
        f"Reclamacoes: {summary['kpis']['total']}",
        f"Procedentes: {summary['kpis']['procedentes']} ({round(summary['kpis']['procedenciaPct'] * 100, 1)}%)",
        f"SLA SAC: {round(productivity['kpis']['slaPct'] * 100, 1)}%",
        f"Sucesso operacional: {round(operation['kpis']['successPct'] * 100, 1)}%",
    ])
    add_bullets("Reclamacoes", [f"{row['label']}: {row['value']}" for row in summary["charts"]["motivo"][:7]])
    add_bullets("Produtividade", [f"{row['analista']}: {row['fechadas']} fechadas, SLA {row['taxa_sla']}%" for row in productivity["tables"]["slaByAnalyst"][:7]])
    add_bullets("Operacao", [f"{row['label']}: sucesso {row['taxa_sucesso']}%, perda {row['taxa_perda']}%" for row in operation["tables"]["hubPerformance"][:7]])
    add_bullets("Mapa", [f"Cidades no mapa: {len(map_payload['points'])}", f"Cidades sem coordenadas: {map_payload['missingCoordinates']}"])
    add_bullets("Alertas", [f"{row['severity']} - {row['area']}: {row['message']}" for row in alerts[:7]])
    add_bullets("Recomendacoes", insights[:7] or ["Manter acompanhamento operacional."])
    prs.save(output)
    return output.getvalue()


def api_filters():
    init_db()
    with connect() as conn:
        def values(field):
            return [
                row["value"]
                for row in conn.execute(
                    f"SELECT DISTINCT {field} AS value FROM complaints WHERE {field} IS NOT NULL ORDER BY value"
                )
            ]
        def union_values(complaint_field, operation_field=None):
            operation_field = operation_field or complaint_field
            rows = conn.execute(
                f"""
                SELECT value FROM (
                    SELECT DISTINCT {complaint_field} AS value FROM complaints WHERE {complaint_field} IS NOT NULL
                    UNION
                    SELECT DISTINCT {operation_field} AS value FROM operations WHERE {operation_field} IS NOT NULL
                )
                ORDER BY value
                """
            ).fetchall()
            return [row["value"] for row in rows]

        return {
            "regiao": union_values("regiao", "regional"),
            "hub": union_values("hub"),
            "cliente": union_values("cliente"),
            "analista": values("analista_sac"),
            "months": union_values("mes"),
        }


def export_csv(query):
    init_db()
    params = {key: values[0] for key, values in query.items() if values and values[0]}
    with connect() as conn:
        max_date = conn.execute("SELECT MAX(data_reclamacao) AS value FROM complaints").fetchone()["value"]
        params["_max_date"] = max_date
        where, args = build_where(params, procedentes_only=params.get("scope") == "procedentes")
        rows = conn.execute(
            f"""
            SELECT
              data_reclamacao, dt_fechamento, mes, mes_fechamento, cliente, cidade, uf, regiao, hub,
              servico, status_reclamacao, validacao, motivo, grupo_motivo, prestador, analista_sac,
              procedente, dias_uteis_fechamento, sla_ok
            FROM complaints
            {where}
            ORDER BY data_reclamacao, analista_sac, cliente
            """,
            args,
        ).fetchall()
    output = io.StringIO()
    writer = csv.writer(output, delimiter=";")
    headers = rows[0].keys() if rows else [
        "data_reclamacao", "dt_fechamento", "mes", "mes_fechamento", "cliente", "cidade",
        "uf", "regiao", "hub", "servico", "status_reclamacao", "validacao", "motivo",
        "grupo_motivo", "prestador", "analista_sac", "procedente", "dias_uteis_fechamento", "sla_ok"
    ]
    writer.writerow(headers)
    for row in rows:
        writer.writerow([row[key] for key in headers])
    return output.getvalue().encode("utf-8-sig")


def parse_multipart_form(headers, rfile):
    content_type = headers.get("Content-Type", "")
    if "multipart/form-data" not in content_type:
        raise ValueError("Envio invalido. Use o formulario de upload do dashboard.")
    try:
        content_length = int(headers.get("Content-Length", "0") or 0)
    except ValueError as exc:
        raise ValueError("Tamanho do arquivo invalido.") from exc
    if content_length <= 0:
        raise ValueError("Arquivo vazio ou nao recebido.")

    body = rfile.read(content_length)
    raw_message = (
        f"Content-Type: {content_type}\r\n"
        "MIME-Version: 1.0\r\n\r\n"
    ).encode("utf-8") + body
    message = BytesParser(policy=default).parsebytes(raw_message)
    fields = {}
    files = {}
    for part in message.iter_parts():
        if part.get_content_disposition() != "form-data":
            continue
        name = part.get_param("name", header="content-disposition")
        if not name:
            continue
        payload = part.get_payload(decode=True) or b""
        filename = part.get_filename()
        if filename:
            files[name] = {"filename": filename, "content": payload}
        else:
            charset = part.get_content_charset() or "utf-8"
            fields[name] = payload.decode(charset, errors="replace")
    return fields, files


class Handler(BaseHTTPRequestHandler):
    def log_message(self, format, *args):
        return

    def send_json(self, payload, status=HTTPStatus.OK):
        data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Cache-Control", "no-store, no-cache, must-revalidate, max-age=0")
        self.send_header("Pragma", "no-cache")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def send_file(self, path, content_type):
        data = path.read_bytes()
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", content_type)
        self.send_header("Cache-Control", "no-store, no-cache, must-revalidate, max-age=0")
        self.send_header("Pragma", "no-cache")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def send_download(self, data, filename, content_type):
        self.send_response(HTTPStatus.OK)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Disposition", f'attachment; filename="{filename}"')
        self.send_header("Cache-Control", "no-store, no-cache, must-revalidate, max-age=0")
        self.send_header("Pragma", "no-cache")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def send_static(self, parsed_path):
        rel = parsed_path.lstrip("/")
        if rel.startswith("static/"):
            rel = rel[len("static/"):]
        target = (STATIC_DIR / rel).resolve()
        if not str(target).startswith(str(STATIC_DIR.resolve())) or not target.is_file():
            return self.send_error(HTTPStatus.NOT_FOUND)
        content_type = "application/octet-stream"
        if target.suffix == ".js":
            content_type = "application/javascript; charset=utf-8"
        elif target.suffix == ".json":
            content_type = "application/json; charset=utf-8"
        elif target.suffix == ".css":
            content_type = "text/css; charset=utf-8"
        elif target.suffix == ".html":
            content_type = "text/html; charset=utf-8"
        return self.send_file(target, content_type)

    def scoped_query(self, parsed):
        from auth import AUTH_ENABLED, require_user
        query = parse_qs(parsed.query)
        if not AUTH_ENABLED:
            return query
        user = require_user(self.headers)
        if user.get("perfil") == "coordenador" and user.get("regional"):
            query["regiao"] = [user["regional"]]
        elif user.get("perfil") == "analista" and user.get("analista_sac"):
            query["analista"] = [user["analista_sac"]]
        return query

    def require_role(self, *roles):
        from auth import AUTH_ENABLED, require_user
        if not AUTH_ENABLED:
            return {"perfil": "admin"}
        user = require_user(self.headers)
        if roles and user.get("perfil") not in roles:
            raise PermissionError("Usuario sem permissao para esta acao.")
        return user

    def do_GET(self):
        parsed = urlparse(self.path)
        from cache_manager import get_cached
        try:
            if parsed.path in ("/", "/index.html"):
                return self.send_file(STATIC_DIR / "index.html", "text/html; charset=utf-8")
            if parsed.path.startswith("/data/") or parsed.path.startswith("/static/"):
                return self.send_static(parsed.path)
            if parsed.path == "/api/filters":
                self.require_role("admin", "coordenador", "analista", "consulta")
                return self.send_json(api_filters())
            if parsed.path == "/api/config":
                self.require_role("admin")
                return self.send_json({"config": load_config(), "path": str(CONFIG_PATH)})
            if parsed.path == "/api/me":
                from auth import current_user, AUTH_ENABLED
                user = current_user(self.headers)
                return self.send_json({"authenticated": bool(user), "authEnabled": AUTH_ENABLED, "user": user})
            query = self.scoped_query(parsed)
        except PermissionError as exc:
            return self.send_json({"ok": False, "error": str(exc)}, HTTPStatus.UNAUTHORIZED)
        if parsed.path in ("/api/summary", "/api/reclamacoes"):
            return self.send_json(get_cached("summary", query, lambda: api_summary(query)))
        if parsed.path in ("/api/productivity", "/api/produtividade"):
            return self.send_json(get_cached("productivity", query, lambda: api_productivity(query)))
        if parsed.path in ("/api/operation", "/api/operacao"):
            return self.send_json(get_cached("operation", query, lambda: api_operation(query)))
        if parsed.path == "/api/kpis":
            return self.send_json(get_cached("kpis", query, lambda: {
                "reclamacoes": api_summary(query)["kpis"],
                "produtividade": api_productivity(query)["kpis"],
                "operacao": api_operation(query)["kpis"],
            }))
        if parsed.path == "/api/prestadores":
            return self.send_json(get_cached("prestadores", query, lambda: api_providers(query)))
        if parsed.path.startswith("/api/score/"):
            dimension = parsed.path.rsplit("/", 1)[-1]
            return self.send_json(get_cached(f"score:{dimension}", query, lambda: api_score_dimension(query, dimension)))
        if parsed.path == "/api/alertas":
            return self.send_json(get_cached("alertas", query, lambda: api_alerts(query)))
        if parsed.path == "/api/insights":
            return self.send_json(get_cached("insights", query, lambda: api_insights(query)))
        if parsed.path == "/api/mapa":
            return self.send_json(get_cached("mapa", query, lambda: api_map(query)))
        if parsed.path == "/api/mapa/status":
            return self.send_json(get_cached("mapa-status", query, api_geo_coverage_snapshot))
        if parsed.path == "/api/geocidades/pendentes":
            return self.send_json(get_cached("geocidades-pendentes", query, api_geocidades_pendentes))
        if parsed.path == "/api/geocidades/pendentes.csv":
            return self.send_download(export_geocidades_pendentes_csv(), "geocidades_pendentes.csv", "text/csv; charset=utf-8")
        if parsed.path == "/api/import/auto/status":
            from auto_import import auto_import_status
            return self.send_json(auto_import_status())
        if parsed.path in ("/api/clientes", "/api/hubs", "/api/regionais"):
            summary = api_summary(query)
            key = {"/api/clientes": "rateClient", "/api/hubs": "rateHub", "/api/regionais": "rateRegion"}[parsed.path]
            return self.send_json({"rows": summary["tables"][key]})
        if parsed.path == "/api/export":
            return self.send_download(export_csv(query), "dashboard_reclamacoes_filtrado.csv", "text/csv; charset=utf-8")
        if parsed.path == "/api/export/excel":
            return self.send_download(export_excel_workbook(query), "dashboard_mms.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        if parsed.path == "/api/export/pdf":
            return self.send_download(export_executive_pdf(query), "resumo_executivo_mms.pdf", "application/pdf")
        if parsed.path in ("/api/export/ppt", "/api/export/pptx"):
            return self.send_download(export_executive_pptx(query), "slides_executivos_mms.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        if parsed.path == "/api/health":
            from cache_manager import cache_available
            from auto_import import auto_import_status
            return self.send_json({
                "ok": True,
                "db": str(DB_PATH),
                "missingDependencies": missing_dependencies(),
                "cacheAvailable": cache_available(),
                "autoImport": auto_import_status(),
            })
        self.send_error(HTTPStatus.NOT_FOUND)

    def do_POST(self):
        parsed = urlparse(self.path)
        if parsed.path == "/api/login":
            try:
                content_length = int(self.headers.get("Content-Length", "0") or 0)
                payload = json.loads(self.rfile.read(content_length).decode("utf-8") or "{}")
                from auth import authenticate
                with connect() as conn:
                    result = authenticate(conn, payload.get("email", ""), payload.get("senha", ""))
                if not result:
                    return self.send_json({"ok": False, "error": "Email ou senha invalidos."}, HTTPStatus.UNAUTHORIZED)
                token, user = result
                data = json.dumps({"ok": True, "user": user}, ensure_ascii=False).encode("utf-8")
                self.send_response(HTTPStatus.OK)
                self.send_header("Content-Type", "application/json; charset=utf-8")
                self.send_header("Set-Cookie", f"mms_session={token}; HttpOnly; SameSite=Lax; Path=/")
                self.send_header("Content-Length", str(len(data)))
                self.end_headers()
                self.wfile.write(data)
                return
            except Exception as exc:
                return self.send_json({"ok": False, "error": str(exc)}, HTTPStatus.BAD_REQUEST)
        if parsed.path == "/api/logout":
            from auth import logout, token_from_headers
            logout(token_from_headers(self.headers))
            data = json.dumps({"ok": True}, ensure_ascii=False).encode("utf-8")
            self.send_response(HTTPStatus.OK)
            self.send_header("Content-Type", "application/json; charset=utf-8")
            self.send_header("Set-Cookie", "mms_session=; Max-Age=0; HttpOnly; SameSite=Lax; Path=/")
            self.send_header("Content-Length", str(len(data)))
            self.end_headers()
            self.wfile.write(data)
            return
        if parsed.path == "/api/config":
            try:
                self.require_role("admin")
                content_length = int(self.headers.get("Content-Length", "0") or 0)
                payload = json.loads(self.rfile.read(content_length).decode("utf-8") or "{}")
                current = load_config()
                updates = payload.get("config", payload)
                for key, value in updates.items():
                    if isinstance(value, dict) and isinstance(current.get(key), dict):
                        current[key].update(value)
                    else:
                        current[key] = value
                save_config(current)
                from cache_manager import clear_api_cache
                clear_api_cache()
                return self.send_json({"ok": True, "config": refresh_config()})
            except Exception as exc:
                return self.send_json({"ok": False, "error": str(exc)}, HTTPStatus.BAD_REQUEST)
        if parsed.path == "/api/geocidades/atualizar":
            try:
                self.require_role("admin", "coordenador")
                try:
                    from geocidades_importer import atualizar_geocidades as atualizar_geocidades_completa
                    result = atualizar_geocidades_completa()
                except FileNotFoundError as exc:
                    return self.send_json({"ok": False, "error": str(exc)}, HTTPStatus.BAD_REQUEST)
                from cache_manager import clear_api_cache
                clear_api_cache()
                return self.send_json({"ok": True, **result})
            except Exception as exc:
                return self.send_json({"ok": False, "error": str(exc)}, HTTPStatus.BAD_REQUEST)
        if parsed.path == "/api/geocidades/salvar":
            try:
                self.require_role("admin", "coordenador")
                content_length = int(self.headers.get("Content-Length", "0") or 0)
                payload = json.loads(self.rfile.read(content_length).decode("utf-8") or "{}")
                return self.send_json(api_save_geocidade(payload))
            except Exception as exc:
                return self.send_json({"ok": False, "error": str(exc)}, HTTPStatus.BAD_REQUEST)
        if parsed.path == "/api/geocidades/ignorar":
            try:
                self.require_role("admin", "coordenador")
                content_length = int(self.headers.get("Content-Length", "0") or 0)
                payload = json.loads(self.rfile.read(content_length).decode("utf-8") or "{}")
                return self.send_json(api_ignore_geocidade(payload))
            except Exception as exc:
                return self.send_json({"ok": False, "error": str(exc)}, HTTPStatus.BAD_REQUEST)
        if parsed.path != "/api/upload":
            return self.send_error(HTTPStatus.NOT_FOUND)
        try:
            self.require_role("admin", "coordenador")
        except PermissionError as exc:
            return self.send_json({"ok": False, "error": str(exc)}, HTTPStatus.UNAUTHORIZED)

        try:
            fields, files = parse_multipart_form(self.headers, self.rfile)
        except ValueError as exc:
            return self.send_json({"ok": False, "error": str(exc)}, HTTPStatus.BAD_REQUEST)

        if "file" not in files:
            return self.send_json({"ok": False, "error": "Nenhum arquivo enviado."}, HTTPStatus.BAD_REQUEST)

        item = files["file"]
        original_name = Path(item["filename"] or "relatorio.xlsx").name
        if not original_name.lower().endswith((".xlsx", ".xls")):
            return self.send_json({"ok": False, "error": "Envie uma planilha .xlsx ou .xls."}, HTTPStatus.BAD_REQUEST)

        mode = "append" if fields.get("mode") == "append" else "replace"
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        saved = UPLOAD_DIR / f"{stamp}_{original_name}"
        saved.write_bytes(item["content"])

        try:
            result = import_workbook(saved, mode=mode)
            from cache_manager import clear_api_cache
            clear_api_cache()
        except Exception as exc:
            return self.send_json({"ok": False, "error": str(exc)}, HTTPStatus.BAD_REQUEST)
        return self.send_json({"ok": True, "result": result})


def run(host, port):
    init_db()
    from auto_import import start_scheduler
    start_scheduler(minutes=30)
    server = ThreadingHTTPServer((host, port), Handler)
    print(f"Dashboard disponível em http://{host}:{port}")
    server.serve_forever()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--host", default="127.0.0.1")
    parser.add_argument("--port", type=int, default=8787)
    parser.add_argument("--import-file")
    parser.add_argument("--mode", choices=["replace", "append"], default="replace")
    parser.add_argument("--serve", action="store_true")
    args = parser.parse_args()

    init_db()
    if args.import_file:
        print(json.dumps(import_workbook(args.import_file, mode=args.mode), ensure_ascii=False, indent=2))
    if args.serve:
        run(args.host, args.port)


if __name__ == "__main__":
    main()
